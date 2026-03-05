# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx"]
# ///
"""Generate Claude_session2.pptx matching session1 design."""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Design tokens ──────────────────────────────────────────
BG_MAIN   = RGBColor(0x1B, 0x19, 0x18)
BG_SECTION = RGBColor(0x23, 0x1F, 0x1D)
CARD_BG   = RGBColor(0x2D, 0x29, 0x26)
CARD_BG2  = RGBColor(0x3A, 0x36, 0x32)
CARD_BG3  = RGBColor(0x30, 0x2D, 0x29)
CARD_BG4  = RGBColor(0x2C, 0x2A, 0x26)

ACCENT    = RGBColor(0xD9, 0x77, 0x57)  # terracotta
ACCENT2   = RGBColor(0x6A, 0x9B, 0xCC)  # blue
ACCENT3   = RGBColor(0x7B, 0xC9, 0xA0)  # green
RED       = RGBColor(0xE0, 0x70, 0x70)

WHITE     = RGBColor(0xFA, 0xF9, 0xF5)
BODY      = RGBColor(0xE0, 0xDE, 0xD8)
MUTED     = RGBColor(0xB0, 0xAE, 0xA5)
MUTED2    = RGBColor(0xC5, 0xC3, 0xBB)
DARK_MUTED = RGBColor(0x6B, 0x69, 0x63)
STRIPE_C  = RGBColor(0x2D, 0x29, 0x26)
LINE_C    = RGBColor(0x4A, 0x47, 0x42)

FONT_JP = "Hiragino Sans"
FONT_EN = "Helvetica Neue"
FONT_CODE = "SF Mono"

SLIDE_W = 12192000
SLIDE_H = 6858000
TOTAL_SLIDES = 19

# ── Helpers ────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill=None, name=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if name:
        shape.name = name
    return shape


def add_rounded_rect(slide, left, top, width, height, fill=None, name=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if name:
        shape.name = name
    # Reduce corner radius
    shape.adjustments[0] = 0.04
    return shape


def _set_font(run, font_name, size, bold=False, color=WHITE):
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font=FONT_JP,
                size=Pt(16), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                name=None, line_spacing=None):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.fill.background()
    if name:
        txbox.name = name
    tf = txbox.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        _set_font(run, font, size, bold, color)
    return txbox


def add_part_label(slide, text):
    return add_textbox(slide, Emu(731520), Emu(365760), Emu(8000000), Emu(274320),
                       text, font=FONT_EN, size=Emu(152400), bold=True, color=ACCENT)


def add_slide_title(slide, text, y=Emu(548640)):
    return add_textbox(slide, Emu(731520), y, Emu(10000000), Emu(548640),
                       text, font=FONT_JP, size=Emu(355600), bold=True, color=WHITE)


def add_page_num(slide, num):
    return add_textbox(slide, Emu(SLIDE_W - 1200000), Emu(SLIDE_H - 400000),
                       Emu(900000), Emu(274320),
                       f"{num} / {TOTAL_SLIDES}", font=FONT_EN,
                       size=Emu(139700), color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, text, y=Emu(6200000)):
    # Footer card
    card = add_shape(slide, Emu(731520), y, Emu(10728960), Emu(457200), fill=CARD_BG4)
    accent = add_shape(slide, Emu(731520), y, Emu(54864), Emu(457200), fill=ACCENT)
    tb = add_textbox(slide, Emu(914400), y + Emu(91440), Emu(10400000), Emu(365760),
                     text, font=FONT_JP, size=Emu(177800), bold=True, color=ACCENT)
    return card, accent, tb


def add_divider(slide, left, top, width):
    line = slide.shapes.add_connector(
        1, left, top, left + width, top  # type 1 = straight
    )
    line.line.color.rgb = LINE_C
    line.line.width = Pt(0.75)
    return line


def add_code_block(slide, left, top, width, height, text):
    """Add a code block with dark background."""
    card = add_rounded_rect(slide, left, top, width, height, fill=RGBColor(0x1B, 0x19, 0x18))
    tb = add_textbox(slide, left + Emu(137160), top + Emu(91440),
                     width - Emu(274320), height - Emu(182880),
                     text, font=FONT_CODE, size=Emu(152400), color=BODY)
    return card, tb


def new_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


# ── Slide builders ─────────────────────────────────────────

def build_title_slide(prs):
    """Slide 01: Title."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    # Decorative elements
    ring = slide.shapes.add_shape(MSO_SHAPE.DONUT, Emu(8382000), Emu(-1143000),
                                   Emu(5486400), Emu(5486400))
    ring.line.fill.background()
    ring.fill.background()
    ring.line.color.rgb = RGBColor(0x3A, 0x36, 0x32)
    ring.line.width = Pt(2)
    ring.name = "Deco Ring"

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9144000), Emu(2743200),
                                     Emu(2286000), Emu(2286000))
    circle.line.fill.background()
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.name = "Deco Circle"

    arc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(7620000), Emu(1600200),
                                  Emu(3200400), Emu(3200400))
    arc.line.fill.background()
    arc.fill.solid()
    arc.fill.fore_color.rgb = ACCENT
    arc.rotation = 0
    arc.name = "Deco Arc"

    small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(457200), Emu(5486400),
                                    Emu(914400), Emu(914400))
    small.line.fill.background()
    small.fill.background()
    small.line.color.rgb = RGBColor(0x3A, 0x36, 0x32)
    small.line.width = Pt(2)
    small.name = "Small Circle"

    # Accent bar
    add_shape(slide, Emu(731520), Emu(1371600), Emu(54864), Emu(3200400), fill=ACCENT)

    # Dots
    for i, x in enumerate([1097280, 1371600, 1645920]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x), Emu(1097280),
                                      Emu(91440), Emu(91440))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.name = f"Dot {i+1}"

    # Stripe
    add_shape(slide, Emu(10058400), Emu(5486400), Emu(1828800), Emu(91440), fill=STRIPE_C)

    # Title
    add_textbox(slide, Emu(1097280), Emu(1371600), Emu(6400800), Emu(1554480),
                "Claude Code\nプロジェクトに定着させる",
                font=FONT_JP, size=Emu(660400), bold=True, color=WHITE)

    # Subtitle
    add_textbox(slide, Emu(1097280), Emu(3109920), Emu(6400800), Emu(548640),
                "〜 CLAUDE.md / Skills / Hooks 〜",
                font=FONT_JP, size=Emu(304800), color=ACCENT)

    # Divider
    add_divider(slide, Emu(1097280), Emu(3840480), Emu(2743200))

    # Meta
    add_textbox(slide, Emu(1097280), Emu(4023360), Emu(6400800), Emu(457200),
                "社内勉強会  |  第2回 / 全3回  |  20 min",
                font=FONT_EN, size=Emu(203200), color=MUTED)


def build_roadmap_slide(prs):
    """Slide 02: Roadmap (3 cards)."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_textbox(slide, Emu(736600), Emu(457200), Emu(7620000), Emu(635000),
                "全3回のロードマップ", font=FONT_JP, size=Emu(406400), bold=True, color=WHITE)

    card_w = Emu(3429000)
    card_h = Emu(3937000)
    card_y = Emu(1270000)
    gap = Emu(228600)

    card_data = [
        {"num": "1", "label": "第1回（済）", "title": "理解する",
         "desc": "基本 / コミュニケーション / Plan Mode",
         "accent": MUTED, "card_fill": CARD_BG3, "label_color": MUTED},
        {"num": "2", "label": "第2回（今日）", "title": "プロジェクトに\n定着させる",
         "desc": "コンテキスト・CLAUDE.md・\nスラッシュコマンド & Skills・Hooks",
         "accent": ACCENT, "card_fill": CARD_BG2, "label_color": ACCENT},
        {"num": "3", "label": "第3回", "title": "大規模開発に挑む",
         "desc": "MCP / worktree / サブエージェント",
         "accent": ACCENT3, "card_fill": CARD_BG3, "label_color": MUTED},
    ]

    for i, d in enumerate(card_data):
        x = Emu(736600) + i * (card_w + gap)

        # Card background
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill=d["card_fill"])

        # Top accent bar
        add_shape(slide, x, card_y, card_w, Emu(63500), fill=d["accent"])

        # Number circle
        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Emu(228600), card_y + Emu(304800),
                                            Emu(609600), Emu(609600))
        num_shape.line.fill.background()
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = d["accent"]
        tf = num_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = d["num"]
        _set_font(run, FONT_EN, Emu(254000), True, WHITE)

        # Label
        add_textbox(slide, x + Emu(965200), card_y + Emu(355600), Emu(2286000), Emu(508000),
                    d["label"], font=FONT_JP, size=Emu(203200), bold=True, color=d["label_color"])

        # Title
        add_textbox(slide, x + Emu(228600), card_y + Emu(1143000), Emu(2984500), Emu(1143000),
                    d["title"], font=FONT_JP, size=Emu(279400), bold=True, color=WHITE)

        # Separator
        add_divider(slide, x + Emu(228600), card_y + Emu(2476500), Emu(2984500))

        # Description
        add_textbox(slide, x + Emu(228600), card_y + Emu(2654300), Emu(2984500), Emu(1016000),
                    d["desc"], font=FONT_JP, size=Emu(190500), color=MUTED2)

    # Goal
    goal = add_textbox(slide, Emu(736600), Emu(5500000), Emu(10728960), Emu(457200),
                       "今日のゴール: Claude Codeをプロジェクト固有の「チームメンバー」として設定できるようになる",
                       font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    add_page_num(slide, 2)


def build_section_slide(prs, part_num, title, subtitle, slide_num):
    """Section divider slide."""
    slide = new_slide(prs)
    set_bg(slide, BG_SECTION)

    # Decorative lines
    line_y = Emu(SLIDE_H // 2)
    add_divider(slide, Emu(731520), line_y, Emu(1828800))
    add_divider(slide, Emu(SLIDE_W - 731520 - 1828800), line_y, Emu(1828800))

    # PART label
    add_textbox(slide, Emu(0), Emu(2000000), Emu(SLIDE_W), Emu(457200),
                f"PART {part_num}", font=FONT_EN, size=Emu(228600), bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Emu(0), Emu(2600000), Emu(SLIDE_W), Emu(762000),
                title, font=FONT_JP, size=Emu(508000), bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Emu(0), Emu(3500000), Emu(SLIDE_W), Emu(457200),
                subtitle, font=FONT_JP, size=Emu(304800),
                color=MUTED2, align=PP_ALIGN.CENTER)

    add_page_num(slide, slide_num)


def build_context_window_slide(prs):
    """Slide 04: What Claude Code sees."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — コンテキストとCLAUDE.md")
    add_slide_title(slide, "Claude Codeは何を見ているか")

    # Hero definition
    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10728960), Emu(548640),
                "コンテキストウィンドウ = Claudeの作業メモリ",
                font=FONT_JP, size=Emu(254000), bold=True, color=MUTED)

    add_divider(slide, Emu(731520), Emu(1800000), Emu(10728960))

    # Memory bar visualization
    bar_y = Emu(2000000)
    bar_h = Emu(548640)
    segments = [
        ("CLAUDE.md", Emu(1200000), ACCENT),
        ("会話履歴", Emu(1600000), ACCENT2),
        ("読んだファイル", Emu(1800000), ACCENT3),
        ("ツール結果", Emu(1400000), RGBColor(0xCC, 0xAA, 0x55)),
        ("空き", Emu(4728960), CARD_BG2),
    ]
    x = Emu(731520)
    for label, w, color in segments:
        seg = add_rounded_rect(slide, x, bar_y, w, bar_h, fill=color)
        seg.adjustments[0] = 0.02
        add_textbox(slide, x + Emu(45720), bar_y + Emu(137160), w - Emu(91440), Emu(274320),
                    label, font=FONT_JP, size=Emu(152400), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        x += w

    # Scale label
    add_textbox(slide, Emu(731520), bar_y + bar_h + Emu(91440), Emu(10728960), Emu(274320),
                "約200,000トークン",
                font=FONT_EN, size=Emu(152400), color=MUTED, align=PP_ALIGN.CENTER)

    # Bullet points
    bullets = [
        "セッション開始時にCLAUDE.mdが自動で読み込まれる",
        "会話が進むにつれ、コンテキストが蓄積していく",
        "ファイルの読み取り、コマンドの実行結果もコンテキストに入る",
    ]
    y = Emu(3200000)
    for b in bullets:
        add_textbox(slide, Emu(914400), y, Emu(10000000), Emu(365760),
                    f"•  {b}", font=FONT_JP, size=Emu(203200), color=BODY)
        y += Emu(400000)

    add_footer(slide, "→ コンテキストウィンドウ = 「Claude Codeが今覚えていること」の総量")
    add_page_num(slide, 4)


def build_compact_slide(prs):
    """Slide 05: Why Claude 'forgets'."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — コンテキストとCLAUDE.md")
    add_slide_title(slide, "なぜ「忘れる」のか")

    # Before/After
    add_textbox(slide, Emu(731520), Emu(1200000), Emu(5000000), Emu(320000),
                "コンパクト: コンテキストの自動圧縮", font=FONT_JP,
                size=Emu(228600), bold=True, color=ACCENT)

    # Before
    add_textbox(slide, Emu(731520), Emu(1600000), Emu(2000000), Emu(274320),
                "Before", font=FONT_EN, size=Emu(177800), bold=True, color=MUTED)
    before_items = ["CLAUDE.md", "指示A", "ファイル読み", "指示B", "実装", "テスト", "修正..."]
    bx = Emu(731520)
    for item in before_items:
        w = Emu(len(item) * 180000 + 200000)
        add_rounded_rect(slide, bx, Emu(1900000), w, Emu(320000), fill=CARD_BG)
        add_textbox(slide, bx + Emu(45720), Emu(1940000), w - Emu(91440), Emu(274320),
                    item, font=FONT_JP, size=Emu(139700), color=BODY, align=PP_ALIGN.CENTER)
        bx += w + Emu(45720)

    # Arrow
    add_textbox(slide, Emu(5000000), Emu(2350000), Emu(600000), Emu(320000),
                "▼", font=FONT_EN, size=Emu(228600), color=ACCENT, align=PP_ALIGN.CENTER)

    # After
    add_textbox(slide, Emu(731520), Emu(2700000), Emu(2000000), Emu(274320),
                "After（コンパクト発動後）", font=FONT_EN, size=Emu(177800), bold=True, color=MUTED)
    after_items = [("CLAUDE.md", ACCENT), ("要約: 古い会話", RED), ("最近のやりとり", ACCENT3)]
    ax = Emu(731520)
    for item, color in after_items:
        w = Emu(len(item) * 180000 + 300000)
        add_rounded_rect(slide, ax, Emu(3000000), w, Emu(320000), fill=CARD_BG)
        add_textbox(slide, ax + Emu(45720), Emu(3040000), w - Emu(91440), Emu(274320),
                    item, font=FONT_JP, size=Emu(139700), color=color, align=PP_ALIGN.CENTER)
        ax += w + Emu(45720)

    add_divider(slide, Emu(731520), Emu(3550000), Emu(10728960))

    # Two columns: 起きること / 対策
    # Left: 起きること
    add_textbox(slide, Emu(731520), Emu(3700000), Emu(4500000), Emu(320000),
                "起きること", font=FONT_JP, size=Emu(203200), bold=True, color=RED)
    problems = [
        "最初の指示や制約が要約・削除される場合がある",
        "「このファイルは触らないで」と言ったのに忘れられる",
        "長いセッションほどリスクが高い",
    ]
    y = Emu(4050000)
    for p in problems:
        add_textbox(slide, Emu(914400), y, Emu(4500000), Emu(274320),
                    f"•  {p}", font=FONT_JP, size=Emu(177800), color=BODY)
        y += Emu(320000)

    # Right: 対策
    add_textbox(slide, Emu(6000000), Emu(3700000), Emu(5000000), Emu(320000),
                "対策", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT3)
    solutions = [
        "重要な情報はCLAUDE.mdに書く（消えない）",
        "作業が変わるタイミングで /clear する",
        "長い作業はタスクを分割する",
    ]
    y = Emu(4050000)
    for s in solutions:
        add_textbox(slide, Emu(6200000), y, Emu(5000000), Emu(274320),
                    f"•  {s}", font=FONT_JP, size=Emu(177800), color=BODY)
        y += Emu(320000)

    add_footer(slide, "→ CLAUDE.mdは毎セッション読み込まれる「永続メモリ」として機能する")
    add_page_num(slide, 5)


def build_claudemd_intro_slide(prs):
    """Slide 06: What is CLAUDE.md."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — コンテキストとCLAUDE.md")
    add_slide_title(slide, "CLAUDE.mdとは")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "人間のオンボーディング資料と同じ発想",
                font=FONT_JP, size=Emu(228600), bold=True, color=ACCENT)

    # Onboarding items
    items = [
        "このプロジェクトはこういう構成です",
        "ビルドはこのコマンドで",
        "こういうルールがあるので守ってください",
        "ここは触らないでください",
    ]
    y = Emu(1650000)
    for item in items:
        add_textbox(slide, Emu(914400), y, Emu(5000000), Emu(274320),
                    f"•  {item}", font=FONT_JP, size=Emu(203200), color=BODY)
        y += Emu(340000)

    add_textbox(slide, Emu(914400), y + Emu(91440), Emu(5000000), Emu(274320),
                "→ これをそのままCLAUDE.mdに書く",
                font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    add_divider(slide, Emu(731520), Emu(3200000), Emu(10728960))

    # Table: 配置場所
    add_textbox(slide, Emu(731520), Emu(3350000), Emu(5000000), Emu(365760),
                "配置場所と読み込みルール", font=FONT_JP,
                size=Emu(228600), bold=True, color=WHITE)

    # Table headers
    col_w = [Emu(3600000), Emu(2400000), Emu(4728960)]
    col_x = [Emu(731520)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    headers = ["場所", "読み込み", "用途"]
    header_y = Emu(3800000)
    row_h = Emu(400000)

    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_shape(slide, cx, header_y, cw, row_h, fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), header_y + Emu(91440), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(177800), bold=True, color=ACCENT)

    rows = [
        ["~/.claude/CLAUDE.md", "常に", "個人の全プロジェクト共通設定"],
        ["プロジェクト/CLAUDE.md", "常に", "プロジェクト固有の設定（チーム共有）"],
        ["サブディレクトリ/CLAUDE.md", "作業時", "モジュール固有の注意事項"],
    ]
    for i, row in enumerate(rows):
        ry = header_y + row_h + Emu(i * 370000)
        for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(365760), fill=CARD_BG)
            fc = BODY if j > 0 else WHITE
            fsz = Emu(177800) if j == 0 else Emu(177800)
            add_textbox(slide, cx + Emu(91440), ry + Emu(68580), cw - Emu(182880), Emu(274320),
                        cell, font=FONT_CODE if j == 0 else FONT_JP, size=fsz, color=fc)

    add_footer(slide, "→ チームでバージョン管理すれば、全員が同じ「教え方」でClaude Codeを使える")
    add_page_num(slide, 6)


def build_claudemd_content_slide(prs):
    """Slide 07: What to write in CLAUDE.md."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — コンテキストとCLAUDE.md")
    add_slide_title(slide, "CLAUDE.mdに書くべき内容")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "4カテゴリに分けて考える", font=FONT_JP, size=Emu(228600), bold=True, color=MUTED)

    # 4 category cards in 2x2
    categories = [
        ("1", "プロジェクト概要", "## リポジトリ概要\nZennに公開する技術記事を\n管理するリポジトリ。"),
        ("2", "よく使うコマンド", "npm run preview  # プレビュー\nnpm run lint     # チェック\nnpm run lint:fix # 自動修正"),
        ("3", "コーディング規約", "## Python\nPython を使う場合は\n必ず uv を使う（pip は使わない）。"),
        ("4", "やってはいけないこと", "## 禁止事項\n- mainに直接プッシュしない\n- .envファイルを編集しない"),
    ]

    card_w = Emu(5200000)
    card_h = Emu(1700000)
    start_x = Emu(731520)
    start_y = Emu(1650000)
    gap_x = Emu(228600)
    gap_y = Emu(228600)

    for i, (num, title, code) in enumerate(categories):
        col = i % 2
        row = i // 2
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)

        add_rounded_rect(slide, cx, cy, card_w, card_h, fill=CARD_BG)

        # Number + title
        num_s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Emu(137160), cy + Emu(137160),
                                        Emu(365760), Emu(365760))
        num_s.line.fill.background()
        num_s.fill.solid()
        num_s.fill.fore_color.rgb = ACCENT
        tf = num_s.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        _set_font(run, FONT_EN, Emu(203200), True, WHITE)

        add_textbox(slide, cx + Emu(594360), cy + Emu(137160), Emu(4200000), Emu(365760),
                    title, font=FONT_JP, size=Emu(203200), bold=True, color=WHITE)

        # Code area
        add_code_block(slide, cx + Emu(137160), cy + Emu(594360),
                       card_w - Emu(274320), card_h - Emu(731520), code)

    add_footer(slide, "→ PRでよく指摘される項目 = CLAUDE.mdに書くべき内容")
    add_page_num(slide, 7)


def build_claudemd_howto_slide(prs):
    """Slide 08: How to start and grow CLAUDE.md."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — コンテキストとCLAUDE.md")
    add_slide_title(slide, "CLAUDE.mdの始め方と育て方")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "/init で自動生成 → 育てる → 刈り込む",
                font=FONT_JP, size=Emu(228600), bold=True, color=MUTED)

    # 3 steps as cards
    steps = [
        ("Step 1", "/init で自動生成", "プロジェクトの技術スタック、\nディレクトリ構成を自動検出\nCLAUDE.md のたたき台が生成される", ACCENT),
        ("Step 2", "チーム固有の情報を追記", "PRでよく指摘される項目\n触ってはいけないファイル\nチーム固有のコーディング規約", ACCENT2),
        ("Step 3", "肥大化を防ぐ", "目安は500行以下\n詳細な仕様は別ファイルに分離\n不要になった情報は定期的に削除", ACCENT3),
    ]

    card_w = Emu(3429000)
    card_h = Emu(2800000)
    start_x = Emu(731520)
    start_y = Emu(1700000)
    gap = Emu(228600)

    for i, (step, title, desc, color) in enumerate(steps):
        cx = start_x + i * (card_w + gap)

        add_rounded_rect(slide, cx, start_y, card_w, card_h, fill=CARD_BG)
        add_shape(slide, cx, start_y, card_w, Emu(54864), fill=color)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(137160), Emu(1200000), Emu(320000),
                    step, font=FONT_EN, size=Emu(177800), bold=True, color=color)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(457200), card_w - Emu(365760), Emu(457200),
                    title, font=FONT_JP, size=Emu(228600), bold=True, color=WHITE)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(1000000), card_w - Emu(365760), Emu(1600000),
                    desc, font=FONT_JP, size=Emu(177800), color=BODY, line_spacing=6)

    # Tips
    add_textbox(slide, Emu(731520), Emu(4700000), Emu(10728960), Emu(365760),
                "Tips: PART 2で紹介する /insights を使えば、改善提案も受けられる",
                font=FONT_JP, size=Emu(177800), color=MUTED)

    add_page_num(slide, 8)


def build_demo_slide(prs, title, slide_num):
    """Demo slide (minimal text)."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    # Large centered "DEMO" text
    add_textbox(slide, Emu(0), Emu(2200000), Emu(SLIDE_W), Emu(762000),
                "DEMO", font=FONT_EN, size=Emu(762000), bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(slide, Emu(0), Emu(3200000), Emu(SLIDE_W), Emu(457200),
                title, font=FONT_JP, size=Emu(254000),
                color=MUTED, align=PP_ALIGN.CENTER)

    add_page_num(slide, slide_num)


def build_checklist_slide(prs):
    """Slide 10: Context management checklist."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — コンテキストとCLAUDE.md")
    add_slide_title(slide, "コンテキスト管理チェックリスト")

    # Table
    col_w = [Emu(3200000), Emu(7528960)]
    col_x = [Emu(731520), Emu(731520 + 3200000)]

    headers = ["タイミング", "やること"]
    header_y = Emu(1300000)
    row_h = Emu(548640)

    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_shape(slide, cx, header_y, cw, Emu(457200), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(137160), header_y + Emu(91440), cw - Emu(274320), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    rows = [
        ("プロジェクト開始時", "/init でCLAUDE.mdを生成。チーム固有の規約を追記"),
        ("セッション開始時", "新しい作業なら /clear で履歴リセット"),
        ("作業中", "重要な制約はCLAUDE.mdに書く（口頭指示はコンパクトで消える）"),
        ("定期的に", "CLAUDE.mdが肥大化していないか確認。/insights で改善点を確認"),
    ]

    for i, (timing, action) in enumerate(rows):
        ry = header_y + Emu(457200) + Emu(i * 500000)
        for j, (cell, cx, cw) in enumerate(zip([timing, action], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(480000), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fc = WHITE if j == 0 else BODY
            fb = True if j == 0 else False
            add_textbox(slide, cx + Emu(137160), ry + Emu(114300), cw - Emu(274320), Emu(320000),
                        cell, font=FONT_JP, size=Emu(190500), bold=fb, color=fc)

    add_footer(slide, "→ コンテキストを制するものが、Claude Codeを制する")
    add_page_num(slide, 10)


def build_builtin_commands_slide(prs):
    """Slide 12: Built-in slash commands."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 2 — スラッシュコマンド & Skills")
    add_slide_title(slide, "組み込みスラッシュコマンド")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "主要な組み込みコマンド", font=FONT_JP, size=Emu(228600), bold=True, color=MUTED)

    # Commands table
    commands = [
        ("/init", "CLAUDE.md を自動生成", "プロジェクト初期設定"),
        ("/clear", "会話履歴をリセット", "作業の切り替え時"),
        ("/config", "権限や設定の確認・変更", "許可リストの管理"),
        ("/insights", "使い方の改善提案", "定期的な振り返り"),
        ("/compact", "コンテキストを手動で圧縮", "メモリが逼迫したとき"),
        ("/review", "差分のレビュー", "コミット前の確認"),
    ]

    col_w = [Emu(1800000), Emu(4200000), Emu(4728960)]
    col_x = [Emu(731520)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    # Headers
    headers = ["コマンド", "用途", "よく使う場面"]
    hy = Emu(1650000)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(400000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), hy + Emu(91440), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(177800), bold=True, color=ACCENT)

    for i, (cmd, usage, scene) in enumerate(commands):
        ry = hy + Emu(400000) + Emu(i * 380000)
        for j, (cell, cx, cw) in enumerate(zip([cmd, usage, scene], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(365760), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j == 0 else FONT_JP
            fc = ACCENT if j == 0 else BODY
            add_textbox(slide, cx + Emu(91440), ry + Emu(68580), cw - Emu(182880), Emu(274320),
                        cell, font=fn, size=Emu(177800), color=fc)

    # Operation image
    add_code_block(slide, Emu(731520), Emu(4200000), Emu(5000000), Emu(640080),
                   "> /clear\n✓ 会話履歴をクリアしました\n> 新しいタスクの指示...")

    add_textbox(slide, Emu(6000000), Emu(4400000), Emu(5000000), Emu(320000),
                "セッション中に / を入力すると\n候補が表示される",
                font=FONT_JP, size=Emu(190500), color=MUTED)

    add_page_num(slide, 12)


def build_skills_intro_slide(prs):
    """Slide 13: Custom slash commands & Skills."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 2 — スラッシュコマンド & Skills")
    add_slide_title(slide, "カスタムスラッシュコマンド & Skills")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "スキル = 自分で作るスラッシュコマンド",
                font=FONT_JP, size=Emu(228600), bold=True, color=ACCENT)

    # Directory structure
    add_code_block(slide, Emu(731520), Emu(1650000), Emu(5000000), Emu(900000),
                   ".claude/skills/\n├── review/\n│   └── SKILL.md      # /review で呼び出せる\n└── deploy.md          # /deploy で呼び出せる")

    add_divider(slide, Emu(731520), Emu(2750000), Emu(10728960))

    # Comparison table
    add_textbox(slide, Emu(731520), Emu(2850000), Emu(5000000), Emu(365760),
                "CLAUDE.mdとの違い", font=FONT_JP, size=Emu(228600), bold=True, color=WHITE)

    col_w = [Emu(2400000), Emu(4164480), Emu(4164480)]
    col_x = [Emu(731520)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    headers = ["", "CLAUDE.md", "スキル"]
    hy = Emu(3300000)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(365760), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), hy + Emu(68580), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(177800), bold=True, color=ACCENT)

    rows = [
        ("読み込み", "常に自動", "必要なときだけ"),
        ("内容", "プロジェクトの基本ルール", "特定の作業手順・知識"),
        ("目安", "500行以下", "必要に応じて分割"),
        ("例", "コーディング規約、ビルドコマンド", "デプロイ手順、レビュー観点"),
    ]

    for i, (label, md, skill) in enumerate(rows):
        ry = hy + Emu(365760) + Emu(i * 340000)
        for j, (cell, cx, cw) in enumerate(zip([label, md, skill], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(320000), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fc = WHITE if j == 0 else BODY
            fb = j == 0
            add_textbox(slide, cx + Emu(91440), ry + Emu(45720), cw - Emu(182880), Emu(274320),
                        cell, font=FONT_JP, size=Emu(177800), bold=fb, color=fc)

    # Key message
    add_footer(slide, "→ CLAUDE.md = 常識、スキル = 専門知識", y=Emu(5200000))
    add_page_num(slide, 13)


def build_skill_types_slide(prs):
    """Slide 14: Skill types and examples."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 2 — スラッシュコマンド & Skills")
    add_slide_title(slide, "スキルの2タイプと実例")

    # Left: Reference content
    lx = Emu(731520)
    ly = Emu(1200000)
    half_w = Emu(5200000)

    add_textbox(slide, lx, ly, half_w, Emu(365760),
                "参照コンテンツ — 知識を渡す", font=FONT_JP,
                size=Emu(228600), bold=True, color=ACCENT2)

    add_code_block(slide, lx, ly + Emu(457200), half_w, Emu(2286000),
                   "---\nname: api-conventions\ndescription: APIエンドポイント作成時\n             の設計パターン\n---\n\nAPIエンドポイントを書くときは以下に従う。\n- RESTfulな命名規則を使う\n- 一貫したエラーフォーマットを返す\n- リクエストバリデーションを含める")

    add_textbox(slide, lx, ly + Emu(2850000), half_w, Emu(365760),
                "→ 関連する作業を検知したとき、自動で読み込まれる",
                font=FONT_JP, size=Emu(177800), color=ACCENT2)

    # Right: Task content
    rx = Emu(6160000)

    add_textbox(slide, rx, ly, half_w, Emu(365760),
                "タスクコンテンツ — 手順を渡す", font=FONT_JP,
                size=Emu(228600), bold=True, color=ACCENT3)

    add_code_block(slide, rx, ly + Emu(457200), half_w, Emu(2286000),
                   "---\nname: fix-issue\ndescription: GitHubのIssueを修正する\ndisable-model-invocation: true\n---\n\nGitHubのIssue: $ARGUMENTS を分析して修正。\n1. gh issue viewでIssue詳細を取得\n2. 関連ファイルを検索\n3. 修正を実装\n4. テストを書いて実行\n5. コミットしてPR作成")

    add_textbox(slide, rx, ly + Emu(2850000), half_w, Emu(365760),
                "→ /fix-issue 42 のように明示的に呼び出す",
                font=FONT_JP, size=Emu(177800), color=ACCENT3)

    add_page_num(slide, 14)


def build_hooks_intro_slide(prs):
    """Slide 17: What are Hooks."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 3 — Hooks")
    add_slide_title(slide, "Hooksとは")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "イベント駆動の自動化",
                font=FONT_JP, size=Emu(228600), bold=True, color=ACCENT)

    # Flow diagram
    flow_y = Emu(1700000)
    # Flow 1: Edit → PostToolUse → lint
    boxes1 = [
        ("ファイル編集", CARD_BG),
        ("PostToolUse", ACCENT),
        ("lint を自動実行", ACCENT3),
    ]
    fx = Emu(731520)
    for label, color in boxes1:
        bw = Emu(2400000)
        add_rounded_rect(slide, fx, flow_y, bw, Emu(400000), fill=color if color != CARD_BG else CARD_BG)
        add_textbox(slide, fx + Emu(45720), flow_y + Emu(91440), bw - Emu(91440), Emu(274320),
                    label, font=FONT_JP, size=Emu(177800), bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        fx += bw + Emu(91440)
        if fx < Emu(8000000):
            add_textbox(slide, fx - Emu(91440), flow_y + Emu(45720), Emu(91440), Emu(320000),
                        "→", font=FONT_EN, size=Emu(228600), color=MUTED, align=PP_ALIGN.CENTER)

    # Flow 2: Bash → PreToolUse → block
    flow_y2 = flow_y + Emu(548640)
    boxes2 = [
        ("Bashコマンド実行", CARD_BG),
        ("PreToolUse", ACCENT),
        ("危険な操作をブロック", RED),
    ]
    fx = Emu(731520)
    for label, color in boxes2:
        bw = Emu(2400000)
        fill = color if color not in (CARD_BG,) else CARD_BG
        add_rounded_rect(slide, fx, flow_y2, bw, Emu(400000), fill=fill)
        tc = WHITE
        add_textbox(slide, fx + Emu(45720), flow_y2 + Emu(91440), bw - Emu(91440), Emu(274320),
                    label, font=FONT_JP, size=Emu(177800), bold=True,
                    color=tc, align=PP_ALIGN.CENTER)
        fx += bw + Emu(91440)
        if fx < Emu(8000000):
            add_textbox(slide, fx - Emu(91440), flow_y2 + Emu(45720), Emu(91440), Emu(320000),
                        "→", font=FONT_EN, size=Emu(228600), color=MUTED, align=PP_ALIGN.CENTER)

    add_divider(slide, Emu(731520), Emu(3000000), Emu(10728960))

    # Settings location + code
    add_textbox(slide, Emu(731520), Emu(3100000), Emu(5000000), Emu(320000),
                "設定場所: .claude/settings.json",
                font=FONT_JP, size=Emu(203200), bold=True, color=WHITE)

    add_code_block(slide, Emu(731520), Emu(3450000), Emu(5486400), Emu(1371600),
                   '{\n  "hooks": {\n    "PostToolUse": [{\n      "matcher": "Edit|Write",\n      "hooks": [{\n        "type": "command",\n        "command": "npm run lint:fix"\n      }]\n    }]\n  }\n}')

    # Events table (right side)
    add_textbox(slide, Emu(6500000), Emu(3100000), Emu(5000000), Emu(320000),
                "主なイベント", font=FONT_JP, size=Emu(203200), bold=True, color=WHITE)

    events = [
        ("PreToolUse", "ツール実行前", "禁止操作のブロック"),
        ("PostToolUse", "ツール実行後", "lint / formatterの自動実行"),
        ("Notification", "通知時", "外部への通知連携"),
        ("Stop", "応答終了時", "レポート生成"),
    ]

    ey = Emu(3500000)
    for ev_name, timing, usage in events:
        add_rounded_rect(slide, Emu(6500000), ey, Emu(4960000), Emu(320000), fill=CARD_BG)
        add_textbox(slide, Emu(6600000), ey + Emu(45720), Emu(1600000), Emu(274320),
                    ev_name, font=FONT_CODE, size=Emu(152400), bold=True, color=ACCENT)
        add_textbox(slide, Emu(8300000), ey + Emu(45720), Emu(3000000), Emu(274320),
                    f"{timing} → {usage}", font=FONT_JP, size=Emu(139700), color=BODY)
        ey += Emu(365760)

    add_footer(slide, "→ CIのpre-commitフックに近い概念。Claude Codeの操作に品質ゲートを挟む")
    add_page_num(slide, 17)


def build_hooks_examples_slide(prs):
    """Slide 18: Hooks practical examples."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 3 — Hooks")
    add_slide_title(slide, "Hooksの実践例")

    # Example 1
    y = Emu(1100000)
    add_textbox(slide, Emu(731520), y, Emu(5000000), Emu(320000),
                "例1: 編集後にlintを自動修正", font=FONT_JP,
                size=Emu(203200), bold=True, color=ACCENT)
    add_code_block(slide, Emu(731520), y + Emu(365760), Emu(5200000), Emu(548640),
                   '{\n  "matcher": "Edit|Write",\n  "hooks": [{ "type": "command", "command": "npm run lint:fix" }]\n}')
    add_textbox(slide, Emu(6200000), y + Emu(365760), Emu(5000000), Emu(548640),
                "Claudeがファイルを\n編集するたび、\nlint違反が自動修正される",
                font=FONT_JP, size=Emu(177800), color=BODY)

    # Example 2
    y = Emu(2300000)
    add_textbox(slide, Emu(731520), y, Emu(5000000), Emu(320000),
                "例2: 危険なコマンドをブロック（PreToolUse）", font=FONT_JP,
                size=Emu(203200), bold=True, color=RED)
    add_code_block(slide, Emu(731520), y + Emu(365760), Emu(5200000), Emu(548640),
                   '{\n  "matcher": "Bash",\n  "hooks": [{ "type": "command",\n    "command": "grep -q \'rm -rf\' && exit 1; exit 0" }]\n}')
    add_textbox(slide, Emu(6200000), y + Emu(365760), Emu(5000000), Emu(548640),
                "hookコマンドにはツール入力が\nstdinで渡される。\nrm -rf を含む操作をブロック",
                font=FONT_JP, size=Emu(177800), color=BODY)

    # Example 3
    y = Emu(3500000)
    add_textbox(slide, Emu(731520), y, Emu(5000000), Emu(320000),
                "例3: Goファイル編集後に型チェック", font=FONT_JP,
                size=Emu(203200), bold=True, color=ACCENT3)
    add_code_block(slide, Emu(731520), y + Emu(365760), Emu(5200000), Emu(548640),
                   '{\n  "matcher": "Edit|Write",\n  "hooks": [{ "type": "command", "command": "go vet ./..." }]\n}')
    add_textbox(slide, Emu(6200000), y + Emu(365760), Emu(5000000), Emu(548640),
                "言語に合わせた\n品質チェックを\n自動実行",
                font=FONT_JP, size=Emu(177800), color=BODY)

    # Points
    add_divider(slide, Emu(731520), Emu(4600000), Emu(10728960))

    points = [
        ("Hooksが失敗するとClaudeに通知 → Claudeが修正を試みる", ACCENT),
        ("「lint実行して」と毎回指示する必要がなくなる", BODY),
        ("チーム全員が同じチェックを自動実行", BODY),
    ]
    py = Emu(4750000)
    for text, color in points:
        add_textbox(slide, Emu(914400), py, Emu(10000000), Emu(274320),
                    f"•  {text}", font=FONT_JP, size=Emu(177800), color=color)
        py += Emu(320000)

    add_page_num(slide, 18)


def build_summary_slide(prs):
    """Slide 19: Summary."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_slide_title(slide, "まとめ / 次回予告", y=Emu(365760))

    add_textbox(slide, Emu(731520), Emu(914400), Emu(10000000), Emu(365760),
                "3つの持ち帰り", font=FONT_JP, size=Emu(254000), bold=True, color=MUTED)

    # 3 takeaway cards
    takeaways = [
        ("1", "CLAUDE.mdでプロジェクトの「常識」を教える",
         "コンテキストの仕組みを理解し、\n永続的な設定ファイルで\n毎回の説明を不要にする", ACCENT),
        ("2", "Skillsで繰り返しの指示をテンプレート化する",
         "参照コンテンツ（知識）と\nタスクコンテンツ（手順）を\n使い分ける", ACCENT2),
        ("3", "Hooksで品質チェックを自動化する",
         "人間が忘れても、\nlint・テスト・型チェックが\n自動で走る", ACCENT3),
    ]

    card_w = Emu(3429000)
    card_h = Emu(2743200)
    start_x = Emu(731520)
    start_y = Emu(1371600)
    gap = Emu(228600)

    for i, (num, title, desc, color) in enumerate(takeaways):
        cx = start_x + i * (card_w + gap)

        add_rounded_rect(slide, cx, start_y, card_w, card_h, fill=CARD_BG)
        add_shape(slide, cx, start_y, card_w, Emu(54864), fill=color)

        # Number
        num_s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Emu(182880), start_y + Emu(182880),
                                        Emu(548640), Emu(548640))
        num_s.line.fill.background()
        num_s.fill.solid()
        num_s.fill.fore_color.rgb = color
        tf = num_s.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        _set_font(run, FONT_EN, Emu(254000), True, WHITE)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(822960), card_w - Emu(365760), Emu(548640),
                    title, font=FONT_JP, size=Emu(203200), bold=True, color=WHITE)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(1554480), card_w - Emu(365760), Emu(1097280),
                    desc, font=FONT_JP, size=Emu(177800), color=BODY)

    # Memorable phrase
    add_textbox(slide, Emu(731520), Emu(4400000), Emu(10728960), Emu(457200),
                "覚え方:  教える → 呼び出す → 自動化する",
                font=FONT_JP, size=Emu(228600), bold=True, color=ACCENT)

    # Next preview
    add_divider(slide, Emu(731520), Emu(4900000), Emu(10728960))

    add_textbox(slide, Emu(731520), Emu(5050000), Emu(10728960), Emu(457200),
                "次回予告: 第3回「大規模開発に挑む」— MCP / worktree / サブエージェント",
                font=FONT_JP, size=Emu(203200), color=MUTED)

    add_page_num(slide, 19)


# ── Main ───────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 01: Title
    build_title_slide(prs)

    # Slide 02: Roadmap
    build_roadmap_slide(prs)

    # Slide 03: Section — PART 1
    build_section_slide(prs, 1, "コンテキストとCLAUDE.md", "毎回の説明、もう要らない", 3)

    # Slide 04: Context window
    build_context_window_slide(prs)

    # Slide 05: Why it forgets
    build_compact_slide(prs)

    # Slide 06: CLAUDE.md intro
    build_claudemd_intro_slide(prs)

    # Slide 07: CLAUDE.md content
    build_claudemd_content_slide(prs)

    # Slide 08: CLAUDE.md how-to
    build_claudemd_howto_slide(prs)

    # Slide 09: Demo — CLAUDE.md & /init
    build_demo_slide(prs, "CLAUDE.md & /init", 9)

    # Slide 10: Checklist
    build_checklist_slide(prs)

    # Slide 11: Section — PART 2
    build_section_slide(prs, 2, "スラッシュコマンド & Skills", "同じ指示、何度も書かない", 11)

    # Slide 12: Built-in commands
    build_builtin_commands_slide(prs)

    # Slide 13: Skills intro
    build_skills_intro_slide(prs)

    # Slide 14: Skill types
    build_skill_types_slide(prs)

    # Slide 15: Demo — Skills
    build_demo_slide(prs, "Skills", 15)

    # Slide 16: Section — PART 3
    build_section_slide(prs, 3, "Hooks", "人間が忘れても、Hooksが動く", 16)

    # Slide 17: Hooks intro
    build_hooks_intro_slide(prs)

    # Slide 18: Hooks examples
    build_hooks_examples_slide(prs)

    # Slide 19: Summary
    build_summary_slide(prs)

    prs.save("slides/Claude_session2.pptx")
    print("✓ slides/Claude_session2.pptx を生成しました（全19スライド）")


if __name__ == "__main__":
    main()
