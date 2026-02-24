# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx"]
# ///
"""第1回 Claude Code ことはじめ 〜基本と対話の心得〜 スライド生成スクリプト."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens ──────────────────────────────────────────
BG_DARK = RGBColor(0x1B, 0x19, 0x18)       # ダーク背景（暖かい黒）
BG_SECTION = RGBColor(0x23, 0x1F, 0x1D)     # セクション背景
BG_CODE = RGBColor(0x14, 0x12, 0x10)        # コードブロック背景
ACCENT = RGBColor(0xD9, 0x77, 0x57)         # Claude オレンジ
ACCENT_LIGHT = RGBColor(0xE8, 0x95, 0x6F)   # 明るいオレンジ
TEXT_WHITE = RGBColor(0xFA, 0xF9, 0xF5)     # Claude Off-white
TEXT_LIGHT = RGBColor(0xC5, 0xC3, 0xBB)     # 暖グレー
TEXT_DIM = RGBColor(0xB0, 0xAE, 0xA5)       # Claude Mid Gray
HIGHLIGHT = RGBColor(0x6A, 0x9B, 0xCC)      # Claude Blue
GREEN = RGBColor(0x7B, 0xC9, 0xA0)          # 暖かい緑
RED = RGBColor(0xE0, 0x70, 0x70)            # 落ち着いた赤

FONT_TITLE = "Helvetica Neue"
FONT_BODY = "Helvetica Neue"
FONT_CODE = "SF Mono"
FONT_JP = "Hiragino Sans"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 16


def set_slide_bg(slide, color):
    """スライド背景を単色で塗りつぶす."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    """角丸矩形を背景として追加."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_JP, line_spacing=1.5):
    """テキストボックスを追加."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=18,
                       color=TEXT_WHITE, font_name=FONT_JP, line_spacing=1.7,
                       bullet=False):
    """複数行テキストを追加. linesは文字列リストまたは(text, color)タプルのリスト."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, clr = line
        else:
            txt, clr = line, color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "• " if bullet else ""
        p.text = prefix + txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = font_name
        p.space_after = Pt(8)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=14):
    """コードブロック風の矩形+テキストを追加."""
    bg = add_shape_bg(slide, left, top, width, height, BG_CODE, corner_radius=0.03)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(0.2),
        width - Inches(0.6), height - Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = HIGHLIGHT
        p.font.name = FONT_CODE
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * 1.6)
    return bg


def add_speaker_notes(slide, text):
    """スピーカーノートを追加."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_page_number(slide, num):
    """ページ番号を追加."""
    add_text_box(
        slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.5),
        Inches(0.8), Inches(0.3),
        f"{num} / {TOTAL}", font_size=10, color=TEXT_DIM,
        alignment=PP_ALIGN.RIGHT, font_name=FONT_BODY
    )


def add_label(slide, left, top, text, color=ACCENT):
    """ラベルバッジを追加."""
    add_text_box(
        slide, left, top, Inches(3), Inches(0.4),
        text, font_size=12, color=color, bold=True,
        font_name=FONT_BODY
    )


# ── Slide builders ─────────────────────────────────────────

def slide_01_title(prs):
    """タイトルスライド."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    add_text_box(
        slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
        "Claude Code ことはじめ", font_size=48, color=TEXT_WHITE,
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
        "〜基本と対話の心得〜", font_size=28, color=ACCENT,
        alignment=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
        "社内勉強会  |  第1回 / 全3回  |  20 min", font_size=16, color=TEXT_DIM,
        alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
    )

    add_speaker_notes(slide,
        "では始めます。自己紹介は軽く済ませて。\n"
        "Claude Codeの社内勉強会、全3回の第1回です。\n"
        "今日のテーマは「理解する」。基本と対話の心得を扱います。")


def slide_02_series_overview(prs):
    """全3回のロードマップ."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
        "全3回のロードマップ", font_size=32, color=TEXT_WHITE, bold=True
    )

    sessions = [
        ("第1回（今日）", "理解する", "基本 / コミュニケーション / Plan Mode", True),
        ("第2回", "プロジェクトに合わせた\nカスタマイズ", "CLAUDE.md / Skills / Hooks", False),
        ("第3回", "外部連携と並列開発", "MCP / worktree / サブエージェント", False),
    ]

    y = Inches(1.6)
    for title, keyword, detail, is_current in sessions:
        bg_color = RGBColor(0x2D, 0x29, 0x26) if is_current else RGBColor(0x25, 0x22, 0x20)

        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.5), bg_color, 0.02)

        add_text_box(
            slide, Inches(1.2), y + Inches(0.15), Inches(2.5), Inches(0.5),
            title, font_size=18, color=ACCENT_LIGHT if is_current else TEXT_DIM, bold=True
        )
        add_text_box(
            slide, Inches(1.2), y + Inches(0.65), Inches(2.5), Inches(0.5),
            keyword, font_size=24, color=TEXT_WHITE, bold=True
        )
        add_text_box(
            slide, Inches(4.5), y + Inches(0.45), Inches(7), Inches(0.7),
            detail, font_size=16, color=TEXT_LIGHT
        )
        y += Inches(1.8)

    # 今日のゴール
    add_text_box(
        slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
        "今日のゴール: 日常業務の一部をClaude Codeに任せられるようになる",
        font_size=16, color=ACCENT, bold=True
    )

    add_page_number(slide, 2)
    add_speaker_notes(slide,
        "3回シリーズの全体像です。\n"
        "第1回、今日は「理解する」。Claude Codeの基本と使い方を押さえます。\n"
        "第2回は「カスタマイズ」。プロジェクトに合わせた設定の話です。\n"
        "第3回は「外部連携と並列開発」。MCP、worktree、サブエージェントを扱います。\n"
        "今日のゴールは、日常業務の一部をClaude Codeに任せられるようになることです。")


def slide_03_first_step(prs):
    """まずは調査だけ任せてみる."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 1 — Claude Code とは")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "まずは調査だけ任せてみる",
        font_size=32, color=TEXT_WHITE, bold=True
    )

    # よくある痛み
    add_text_box(
        slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
        "よくある痛み: 調査・影響範囲の把握で時間が溶ける",
        font_size=18, color=TEXT_DIM
    )

    # Before / After
    before_y = Inches(2.8)
    add_shape_bg(slide, Inches(0.8), before_y, Inches(5.3), Inches(2.0),
                 RGBColor(0x2D, 0x1E, 0x1E), 0.03)
    add_text_box(
        slide, Inches(1.2), before_y + Inches(0.15), Inches(4.8), Inches(0.35),
        "Before", font_size=16, color=RED, bold=True, font_name=FONT_BODY
    )
    add_text_box(
        slide, Inches(1.2), before_y + Inches(0.6), Inches(4.8), Inches(1.2),
        "チケット1件の調査に30分。\nログを追い、参照箇所を探し、\n影響範囲を整理する",
        font_size=16, color=TEXT_LIGHT
    )

    add_shape_bg(slide, Inches(6.5), before_y, Inches(5.8), Inches(2.0),
                 RGBColor(0x1E, 0x2D, 0x1E), 0.03)
    add_text_box(
        slide, Inches(6.9), before_y + Inches(0.15), Inches(5.0), Inches(0.35),
        "After", font_size=16, color=GREEN, bold=True, font_name=FONT_BODY
    )
    add_text_box(
        slide, Inches(6.9), before_y + Inches(0.6), Inches(5.0), Inches(1.2),
        "Claude Codeに任せると数分で返ってくる。\n調査結果・原因候補・影響範囲が\nまとまる",
        font_size=16, color=TEXT_LIGHT
    )

    # 入口
    add_text_box(
        slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
        "入口: 最初は実装ではなく「調査・要約・差分整理」だけ",
        font_size=18, color=TEXT_LIGHT
    )

    # 下部メッセージ
    add_text_box(
        slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
        "→ いきなり実装させず、業務の一部から始める",
        font_size=18, color=ACCENT, bold=True
    )

    add_page_number(slide, 3)
    add_speaker_notes(slide,
        "調査や影響範囲の把握って、時間が溶けやすい作業です。"
        "ここが最初の狙いどころになります。\n"
        "いきなり実装を任せる必要はありません。"
        "読み取り専用で始められるのでリスクも低い。\n"
        "Before/Afterは体感の話なので、人やリポジトリによって差はあります。\n"
        "ただ「調査の時間が減る」のは、共通して実感しやすいポイントです。\n"
        "この感覚を持ってもらった上で、もう少し詳しく見ていきます。")


def slide_04_what_is_cc(prs):
    """Claude Code でできること."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 1 — Claude Code とは")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "Claude Code でできること", font_size=32, color=TEXT_WHITE, bold=True
    )

    add_multiline_text(
        slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
        [
            "ターミナルで動くAIエージェント\n  指示は自然言語でOK",
            "調査 → 計画 → 実装 → 検証 → 取りまとめ\n  を一連で扱う",
            "デフォルトは読み取り専用\n  ファイル編集やコマンド実行は許可制",
            "コード補完ではなく「作業を進める」ツール",
        ],
        font_size=18, bullet=True, line_spacing=1.6
    )

    # 右側にフロー図
    flow_y = Inches(2.0)
    steps = ["調査", "計画", "実装", "検証", "要約"]
    for i, step in enumerate(steps):
        clr = ACCENT if i == 0 else TEXT_WHITE
        add_shape_bg(
            slide, Inches(8.0), flow_y, Inches(3.0), Inches(0.7),
            RGBColor(0x2D, 0x29, 0x26), 0.05
        )
        add_text_box(
            slide, Inches(8.0), flow_y + Inches(0.1), Inches(3.0), Inches(0.5),
            f"{'→ ' if i > 0 else ''}{step}",
            font_size=20, color=clr, bold=True, alignment=PP_ALIGN.CENTER
        )
        flow_y += Inches(0.9)

    add_page_number(slide, 4)
    add_speaker_notes(slide,
        "Claude Codeは、ターミナルで動くAIエージェントです。\n"
        "「AIにコードを書かせるツールでしょ？」と思われがちですが、"
        "コード生成はできることの一部にすぎません。\n"
        "調査、計画、実装、検証、取りまとめ。"
        "これを一つのセッションで連続的に扱えます。\n"
        "安全面としては、デフォルトは読み取り専用で、ファイル編集やコマンド実行には"
        "都度許可が要ります。勝手に何かを壊す心配はありません。")


def slide_05_tools(prs):
    """Claude Code のツール群."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 1 — Claude Code とは")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "Claude Code のツール群", font_size=32, color=TEXT_WHITE, bold=True
    )

    # 基本ツール（上段・4カード横並び）
    basic_tools = [
        ("Read / Grep / Glob", "ファイルの閲覧\nコード検索"),
        ("Edit / Write", "ファイルの編集\n作成"),
        ("WebSearch /\nWebFetch", "Web検索\nページ取得"),
        ("Task", "サブエージェント\nによる並列探索"),
    ]

    x = Inches(0.8)
    card_w = Inches(2.7)
    for name, desc in basic_tools:
        add_shape_bg(slide, x, Inches(1.9), card_w, Inches(1.4),
                     RGBColor(0x25, 0x22, 0x20), 0.03)
        add_text_box(
            slide, x + Inches(0.2), Inches(2.0), Inches(2.3), Inches(0.5),
            name, font_size=14, color=HIGHLIGHT, bold=True, font_name=FONT_CODE
        )
        add_text_box(
            slide, x + Inches(0.2), Inches(2.55), Inches(2.3), Inches(0.7),
            desc, font_size=13, color=TEXT_LIGHT
        )
        x += Inches(2.9)

    # Bash（下段・大きく）
    add_shape_bg(slide, Inches(0.8), Inches(3.6), Inches(11.5), Inches(3.0),
                 RGBColor(0x1E, 0x2A, 0x35), 0.03)
    add_text_box(
        slide, Inches(1.2), Inches(3.7), Inches(10), Inches(0.5),
        "Bash — あらゆるコマンドを実行できる", font_size=22, color=ACCENT, bold=True
    )

    code_lines = (
        'say "完了しました"           # 音声で通知\n'
        "marp slides.md --pptx        # スライド生成\n"
        "ffmpeg -i video.mp4 out.gif  # 動画をGIFに変換\n"
        "git log --oneline -10        # 直近のコミットを確認\n"
        "npm test                     # テスト実行"
    )
    add_code_block(
        slide, Inches(1.2), Inches(4.3), Inches(10.5), Inches(2.1),
        code_lines, font_size=14
    )

    add_text_box(
        slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
        "→ 開発作業に限らない。コマンドで実現できることは、何でも",
        font_size=16, color=ACCENT
    )

    add_page_number(slide, 5)
    add_speaker_notes(slide,
        "Claude Codeが使えるツールを見てみます。\n"
        "基本ツールとして、ファイルの読み書き・検索、Web検索、サブエージェントがあります。"
        "これらはClaude Codeの内蔵ツールです。\n"
        "注目はBash。文字通り「あらゆるコマンド」を実行できます。\n"
        "テストやビルドはもちろん、marpでスライドを生成する、"
        "ffmpegで動画を変換する。こういったことも全部Bashの守備範囲です。\n"
        "コマンドで実現できることは、原理的にすべてClaude Codeに任せられます。\n"
        "どのツールを使うかは指定不要で、指示に応じてClaude Codeが自動で組み合わせます。\n"
        "ちなみにスライドではgitを例に出していますが、"
        "Perforceなら p4 changes -m 10 で同じようなことができます。"
        "コマンドラインから操作できるツールなら何でも使えます。")


def slide_06_role_comparison(prs):
    """ChatGPT・Copilotとの使い分け."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
        "ChatGPT・Copilotとの使い分け", font_size=32, color=TEXT_WHITE, bold=True
    )

    tools = [
        ("ChatGPT", "設計議論・仕様の言語化", "壁打ち", RGBColor(0x74, 0xAA, 0x9C)),
        ("GitHub Copilot", "IDE内のコード補完", "書く", RGBColor(0x79, 0xB8, 0xFF)),
        ("Claude Code", "リポジトリ全体の作業フロー", "進める", ACCENT),
    ]

    # 「得意領域」ヘッダー
    add_text_box(
        slide, Inches(1.3), Inches(1.4), Inches(5.0), Inches(0.4),
        "得意領域", font_size=14, color=TEXT_DIM, font_name=FONT_BODY
    )

    y = Inches(1.8)
    for name, desc, keyword, color in tools:
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.4), RGBColor(0x25, 0x22, 0x20), 0.02)

        add_text_box(
            slide, Inches(1.3), y + Inches(0.15), Inches(3.0), Inches(0.5),
            name, font_size=22, color=color, bold=True
        )
        add_text_box(
            slide, Inches(1.3), y + Inches(0.7), Inches(5.0), Inches(0.5),
            desc, font_size=16, color=TEXT_LIGHT
        )
        add_text_box(
            slide, Inches(8.5), y + Inches(0.3), Inches(3.5), Inches(0.6),
            keyword, font_size=20, color=color, bold=True,
            alignment=PP_ALIGN.RIGHT
        )
        y += Inches(1.7)

    add_text_box(
        slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
        "→ 排他的ではなく補完的。併用するのがおすすめ",
        font_size=16, color=TEXT_DIM
    )

    add_page_number(slide, 6)
    add_speaker_notes(slide,
        "この3つのツールは競合ではなく、補完的な関係です。得意領域が違います。\n"
        "ChatGPTは壁打ちに強い。設計議論、仕様の言語化、アプローチの比較に向いています。\n"
        "Copilotはコードを書くことに強い。IDEから離れずに入力を補助してくれます。\n"
        "Claude Codeは一連の作業フローを進めることに強い。"
        "調査から実装、差分整理まで通しで扱えます。\n"
        "大事なのは「どれか1つ」ではなく、使い分けることです。")


def slide_07_workflow_shift(prs):
    """開発フローがどう変わるか."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
        "開発フローがどう変わるか", font_size=32, color=TEXT_WHITE, bold=True
    )

    # Column headers
    add_text_box(
        slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.5),
        "従来の開発フロー", font_size=20, color=ACCENT, bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.5),
        "Claude Code を使った開発フロー", font_size=20, color=HIGHLIGHT, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    card_h = Inches(0.6)
    gap = Inches(0.15)
    dev_bg = RGBColor(0x35, 0x2A, 0x22)
    claude_bg = RGBColor(0x1E, 0x2A, 0x35)

    # ── Left column: all Developer ──
    left_steps = [
        "コードを読む",
        "方針を考える",
        "コードを書く",
        "テスト実行",
        "デバッグ・修正",
        "差分整理・PR作成",
    ]

    left_x = Inches(0.5)
    left_w = Inches(5.5)
    y = Inches(1.9)

    for step in left_steps:
        add_shape_bg(slide, left_x, y, left_w, card_h, dev_bg, 0.04)
        add_text_box(
            slide, left_x + Inches(0.3), y + Inches(0.1), Inches(3.5), Inches(0.4),
            step, font_size=15, color=ACCENT_LIGHT, bold=True
        )
        add_text_box(
            slide, left_x + left_w - Inches(2.0), y + Inches(0.1),
            Inches(1.7), Inches(0.4),
            "Developer", font_size=11, color=ACCENT,
            alignment=PP_ALIGN.RIGHT, font_name=FONT_BODY
        )
        y += card_h + gap

    # ── Right column: mixed Claude / Developer ──
    right_steps = [
        ("目標・背景を伝える", True, None),
        ("コード調査", False, None),
        ("計画提案", False, "レビュー"),
        ("実装", False, None),
        ("テスト・検証", False, "確認"),
        ("差分整理・PR作成", False, None),
    ]

    right_x = Inches(6.8)
    right_w = Inches(6.0)
    y = Inches(1.9)

    for text, is_dev, checkpoint in right_steps:
        if checkpoint:
            c_w = Inches(3.4)
            arrow_w = Inches(0.4)
            d_w = Inches(2.2)

            add_shape_bg(slide, right_x, y, c_w, card_h, claude_bg, 0.04)
            add_text_box(
                slide, right_x + Inches(0.3), y + Inches(0.1),
                Inches(2.0), Inches(0.4),
                text, font_size=15, color=HIGHLIGHT, bold=True
            )
            add_text_box(
                slide, right_x + Inches(2.3), y + Inches(0.1),
                Inches(0.8), Inches(0.4),
                "Claude", font_size=11, color=HIGHLIGHT,
                alignment=PP_ALIGN.RIGHT, font_name=FONT_BODY
            )
            add_text_box(
                slide, right_x + c_w, y + Inches(0.1),
                arrow_w, Inches(0.4),
                "→", font_size=15, color=TEXT_DIM,
                alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
            )
            d_x = right_x + c_w + arrow_w
            add_shape_bg(slide, d_x, y, d_w, card_h, dev_bg, 0.04)
            add_text_box(
                slide, d_x + Inches(0.2), y + Inches(0.1),
                Inches(1.8), Inches(0.4),
                checkpoint, font_size=15, color=ACCENT_LIGHT, bold=True,
                alignment=PP_ALIGN.CENTER
            )
        else:
            bg = dev_bg if is_dev else claude_bg
            txt_c = ACCENT_LIGHT if is_dev else HIGHLIGHT
            role = "Developer" if is_dev else "Claude"
            role_c = ACCENT if is_dev else HIGHLIGHT

            add_shape_bg(slide, right_x, y, right_w, card_h, bg, 0.04)
            add_text_box(
                slide, right_x + Inches(0.3), y + Inches(0.1),
                Inches(3.5), Inches(0.4),
                text, font_size=15, color=txt_c, bold=True
            )
            add_text_box(
                slide, right_x + right_w - Inches(2.0), y + Inches(0.1),
                Inches(1.7), Inches(0.4),
                role, font_size=11, color=role_c,
                alignment=PP_ALIGN.RIGHT, font_name=FONT_BODY
            )

        y += card_h + gap

    # Bottom messages
    add_text_box(
        slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4),
        "→ 開発者の出番は「方向づけ」と「判断」の2回だけになる",
        font_size=18, color=ACCENT, bold=True
    )

    add_page_number(slide, 7)
    add_speaker_notes(slide,
        "従来の開発では、全工程を自分で回します。6ステップすべてが自分の作業です。\n"
        "Claude Codeを使うと、開発者の出番は2箇所に絞られます。"
        "最初の方向づけと、途中の判断。\n"
        "これは「サボる」ではなく「判断に集中する」ということです。\n"
        "「コード補完」と「エージェント」の本質的な違いはここにあります。")


def slide_08_task_example(prs):
    """具体例で見てみる."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    dev_bg = RGBColor(0x35, 0x2A, 0x22)
    claude_bg = RGBColor(0x1E, 0x2A, 0x35)

    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
        "具体例で見てみる", font_size=32, color=TEXT_WHITE, bold=True
    )
    add_text_box(
        slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
        "例：「CIが落ちた。原因を調査して修正案を出して」",
        font_size=18, color=TEXT_DIM
    )

    # タイムライン
    line_y = Inches(3.05)
    add_shape_bg(slide, Inches(0.6), line_y, Inches(12.0), Inches(0.04),
                 RGBColor(0x3A, 0x36, 0x32), 0.0)

    cards = [
        {
            "role": "Developer",
            "action": "方向づけ",
            "quote": "\"CIが落ちた。\n ログを要約して\n 原因候補を\n 3つ挙げて\"",
            "is_dev": True,
        },
        {
            "role": "Claude",
            "action": "調査 → 計画提案",
            "quote": "\"失敗箇所はX。\n 依存関係の不整合\n が濃厚。原因候補\n を3つ提示します\"",
            "is_dev": False,
        },
        {
            "role": "Developer",
            "action": "判断",
            "quote": "\"候補2を深掘り\n して。既存仕様は\n 変えないで\"",
            "is_dev": True,
        },
        {
            "role": "Claude",
            "action": "実行 → 取りまとめ",
            "quote": "\"根拠と修正案を\n 整理。テストも\n 追加済み。\"",
            "is_dev": False,
        },
    ]

    card_w = Inches(2.7)
    card_h = Inches(3.5)
    gap = Inches(0.25)
    start_x = Inches(0.6)
    card_y = Inches(1.9)

    for i, card in enumerate(cards):
        x = start_x + (card_w + gap) * i
        bg = dev_bg if card["is_dev"] else claude_bg
        role_color = ACCENT if card["is_dev"] else HIGHLIGHT
        label_color = ACCENT_LIGHT if card["is_dev"] else HIGHLIGHT

        add_shape_bg(slide, x, card_y, card_w, card_h, bg, 0.03)

        add_text_box(
            slide, x + Inches(0.2), card_y + Inches(0.15),
            Inches(2.3), Inches(0.35),
            card["role"], font_size=13, color=role_color, bold=True,
            font_name=FONT_BODY
        )
        add_text_box(
            slide, x + Inches(0.2), card_y + Inches(0.55),
            Inches(2.3), Inches(0.4),
            card["action"], font_size=17, color=label_color, bold=True
        )
        add_text_box(
            slide, x + Inches(0.2), card_y + Inches(1.1),
            Inches(2.3), Inches(2.2),
            card["quote"], font_size=13, color=TEXT_LIGHT,
            font_name=FONT_CODE, line_spacing=1.5
        )

    add_text_box(
        slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
        "→ 「丸投げ」ではなく「要所で判断する」。この感覚が重要",
        font_size=18, color=ACCENT, bold=True
    )

    add_page_number(slide, 8)
    add_speaker_notes(slide,
        "具体的なタスクで流れを見てみます。「CIが落ちた」という日常的なシナリオです。\n"
        "開発者がやるのは、目標を伝えること（30秒）と、計画をレビューすること（1分）。\n"
        "Claudeが調査・計画提案と結果の取りまとめを担当します。\n"
        "開発者の介入は2回。方向づけと判断だけです。\n"
        "「丸投げ」ではありません。要所で判断を入れるから品質を保てます。\n"
        "原因調査に30分かかるところが、5分で俯瞰できるイメージです。\n"
        "ここからPART 2で「伝え方のコツ」を見ていきます。")


def slide_09_section_communication(prs):
    """セクション区切り — コミュニケーション."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SECTION)

    add_text_box(
        slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
        "PART 2", font_size=18, color=ACCENT, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
    )
    add_text_box(
        slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.0),
        "コミュニケーションの心得", font_size=40, color=TEXT_WHITE,
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.0),
        "伝え方で結果が変わるのは、\n人間相手と同じ",
        font_size=24, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER
    )

    add_page_number(slide, 9)
    add_speaker_notes(slide,
        "PART 2、コミュニケーションの話に入ります。\n"
        "先ほどの例で、開発者の発話は2回だけでした。短いですよね。"
        "でも、伝え方で結果は大きく変わります。\n"
        "公式のベストプラクティスにも「Claude Codeへの伝え方は、"
        "結果の質に大きく影響する」と書かれています。\n"
        "人間相手と同じです。何をしてほしいか明確に伝えないと、"
        "期待通りの結果は返ってきません。")


def slide_10_agent_mental_model(prs):
    """チャットボットとエージェントの違い."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 2 — コミュニケーション")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "チャットボットとエージェントの違い", font_size=32, color=TEXT_WHITE, bold=True
    )

    # 左: チャットボット
    add_shape_bg(slide, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5),
                 RGBColor(0x25, 0x22, 0x20), 0.02)
    add_text_box(
        slide, Inches(1.2), Inches(2.2), Inches(4.5), Inches(0.5),
        "チャットボット", font_size=22, color=RED, bold=True
    )
    add_text_box(
        slide, Inches(1.2), Inches(2.7), Inches(4.5), Inches(0.4),
        "電話サポートに近い", font_size=15, color=TEXT_DIM
    )
    add_multiline_text(
        slide, Inches(1.2), Inches(3.3), Inches(4.5), Inches(2.5),
        [
            "質問 → 回答 → 質問 → 回答",
            "1回のやりとりが完結",
            "ユーザーが手順を考え、\n  1ステップずつ指示する",
        ],
        font_size=16, color=TEXT_LIGHT, bullet=True, line_spacing=1.7
    )

    # 右: エージェント
    add_shape_bg(slide, Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.5),
                 RGBColor(0x1E, 0x2A, 0x35), 0.02)
    add_text_box(
        slide, Inches(7.4), Inches(2.2), Inches(4.5), Inches(0.5),
        "エージェント（Claude Code）", font_size=22, color=GREEN, bold=True
    )
    add_text_box(
        slide, Inches(7.4), Inches(2.7), Inches(4.5), Inches(0.4),
        "仕事を任せた同僚に近い", font_size=15, color=TEXT_DIM
    )
    add_multiline_text(
        slide, Inches(7.4), Inches(3.3), Inches(4.5), Inches(2.5),
        [
            "ゴールを伝えたら自律的に実行",
            "ファイル読み書き・コマンド実行・\n  Web検索",
            "調査 → 計画 → 実装を\n  一連で判断・実行する",
        ],
        font_size=16, color=TEXT_LIGHT, bullet=True, line_spacing=1.7
    )

    add_text_box(
        slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
        "→ 「1行ずつ指示する」のではなく「ゴールと背景を伝えて任せる」",
        font_size=16, color=ACCENT
    )

    add_page_number(slide, 10)
    add_speaker_notes(slide,
        "チャットボットとエージェントは根本的に違います。\n"
        "チャットボットは電話サポートに近い。こちらが1ステップずつ指示して、"
        "その都度答えが返ってきます。\n"
        "エージェントは仕事を任せた同僚に近い。ゴールと背景を伝えたら、"
        "やり方は本人が判断します。\n"
        "同僚に仕事を任せるとき、1行ずつ「次はこのファイルを開いて、"
        "次はこの関数を見て」とは指示しませんよね。\n"
        "「CIが落ちてるから原因を調べて」と伝えて、やり方は任せる。\n"
        "この「委任」の感覚が、Claude Codeではとても大事です。")


def slide_11_good_instructions(prs):
    """指示の出し方."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 2 — コミュニケーション")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "指示の出し方", font_size=32, color=TEXT_WHITE, bold=True
    )

    elements = [
        ("目標", "何を達成したいか", ACCENT),
        ("背景・制約", "なぜ必要か / やってはいけないこと", HIGHLIGHT),
        ("現状", "今どうなっているか", GREEN),
        ("期待する出力", "どういう形式で結果がほしいか", RGBColor(0xB8, 0x8A, 0xD0)),
    ]

    y = Inches(1.9)
    for label, desc, color in elements:
        add_shape_bg(slide, Inches(0.8), y, Inches(5.0), Inches(0.9),
                     RGBColor(0x25, 0x22, 0x20), 0.03)
        add_text_box(
            slide, Inches(1.2), y + Inches(0.05), Inches(2.0), Inches(0.45),
            label, font_size=20, color=color, bold=True
        )
        add_text_box(
            slide, Inches(1.2), y + Inches(0.45), Inches(4.2), Inches(0.4),
            desc, font_size=14, color=TEXT_LIGHT
        )
        y += Inches(1.05)

    # Bad
    add_text_box(
        slide, Inches(6.5), Inches(1.9), Inches(3), Inches(0.4),
        "Bad", font_size=16, color=RED, bold=True
    )
    add_code_block(
        slide, Inches(6.5), Inches(2.4), Inches(5.8), Inches(0.7),
        "この機能を改善して", font_size=14
    )

    # Good
    add_text_box(
        slide, Inches(6.5), Inches(3.5), Inches(3), Inches(0.4),
        "Good", font_size=16, color=GREEN, bold=True
    )
    add_code_block(
        slide, Inches(6.5), Inches(4.0), Inches(5.8), Inches(2.6),
        "CIが落ちている。\n"
        "現状：昨日のリファクタ以降失敗するようになった。\n"
        "目標：ログを要約し、原因候補を3つ挙げて。\n"
        "調査結果は箇条書きでまとめて。\n"
        "コードの変更はまだしないこと。",
        font_size=13
    )

    add_page_number(slide, 11)
    add_speaker_notes(slide,
        "指示の4要素。目標、背景・制約、現状、期待する出力です。\n"
        "全部揃っている必要はありません。ただし、目標は必須です。\n"
        "悪い例を見てください。「この機能を改善して」。何を？どう？がまったくわかりません。\n"
        "良い例は、さっきと同じCIのシナリオです。"
        "現状・目標・制約・出力形式が明確になっています。\n"
        "「コードの変更はまだしないこと」が地味に大事で、"
        "やってほしくないことを書くと暴走を防げます。\n"
        "指示が具体的であるほど、結果も具体的になります。")


def slide_12_iteration_loop(prs):
    """たたき台から始める."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 2 — コミュニケーション")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "たたき台から始める", font_size=32, color=TEXT_WHITE, bold=True
    )

    # ループ図（横4ステップ）
    steps_loop = [
        ("1. 指示を出す", ACCENT),
        ("2. 出力を確認", TEXT_WHITE),
        ("3. フィードバック", GREEN),
        ("4. 改善版", TEXT_WHITE),
    ]

    x = Inches(0.8)
    for text, color in steps_loop:
        add_shape_bg(slide, x, Inches(2.0), Inches(2.7), Inches(0.9),
                     RGBColor(0x25, 0x22, 0x20), 0.03)
        add_text_box(
            slide, x + Inches(0.1), Inches(2.15), Inches(2.5), Inches(0.6),
            text, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER
        )
        x += Inches(3.0)

    # 3つのテクニック
    add_text_box(
        slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
        "3つのテクニック", font_size=22, color=TEXT_WHITE, bold=True
    )

    tips = [
        ("「なぜ」を聞く",
         "期待と違ったら、修正指示の前に理由を聞く。\n自分の指示の問題が見えることがある"),
        ("提案形式で伝える",
         "「Xをして」より「Xを考えているが、どう思うか」。\n別の選択肢が出てくる"),
        ("手直ししたら伝える",
         "自分で編集した場合はClaudeに共有する。\n伝えないと古い状態で動き続ける"),
    ]

    y = Inches(4.2)
    for i, (title, desc) in enumerate(tips):
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.1),
                     RGBColor(0x25, 0x22, 0x20), 0.02)
        add_text_box(
            slide, Inches(1.3), y + Inches(0.1), Inches(3.0), Inches(0.5),
            title, font_size=18, color=HIGHLIGHT, bold=True
        )
        add_text_box(
            slide, Inches(4.5), y + Inches(0.1), Inches(7.5), Inches(0.9),
            desc, font_size=14, color=TEXT_LIGHT
        )
        y += Inches(1.2)

    add_text_box(
        slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
        "→ 「全然違う」より「この部分は良い、ここを変えて」",
        font_size=16, color=ACCENT
    )

    add_page_number(slide, 12)
    add_speaker_notes(slide,
        "最初から完璧を求めない。これが一番大事なマインドセットです。\n"
        "まず出力を見て、良い点と悪い点を具体的にフィードバックして、"
        "改善版を出してもらいます。\n"
        "テクニックを3つ紹介します。\n"
        "1つ目：「なぜそうしたのか」と聞く。"
        "自分の指示に問題があったケースも多いです。\n"
        "2つ目：命令ではなく提案で伝える。「Xを考えているけど、どう思う？」と聞くと、"
        "自分では気づかなかった選択肢が出てきます。\n"
        "3つ目：自分でコードを直接編集したら、必ずClaudeに伝えてください。"
        "伝えないと古い状態を前提に動いてしまいます。")


def slide_13_section_plan_mode(prs):
    """セクション区切り — Plan Mode."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SECTION)

    add_text_box(
        slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
        "PART 3", font_size=18, color=ACCENT, bold=True,
        alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
    )
    add_text_box(
        slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.0),
        "Plan Mode", font_size=44, color=TEXT_WHITE,
        bold=True, alignment=PP_ALIGN.CENTER, font_name=FONT_BODY
    )
    add_text_box(
        slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.0),
        "いきなりコードを書かせない",
        font_size=24, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER
    )

    add_page_number(slide, 13)
    add_speaker_notes(slide,
        "PART 3、Plan Modeの話に入ります。\n"
        "「いきなりコードを書かせない」。これが合言葉です。\n"
        "PART 1-2で「何ができるか」と「どう伝えるか」を見てきました。\n"
        "PART 3は「どう進めるか」。ワークフローの話です。\n"
        "計画から始めることで、手戻りを大幅に減らせます。")


def slide_14_plan_mode_what(prs):
    """Plan Mode とは."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 3 — Plan Mode")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "Plan Mode とは", font_size=32, color=TEXT_WHITE, bold=True
    )

    # 説明
    add_multiline_text(
        slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.8),
        [
            "読み取り専用のモード",
            "ファイル閲覧・検索・Web検索は可能",
            "ファイル編集・コマンド実行は不可",
            "調査と計画立案に集中するためのモード",
        ],
        font_size=18, bullet=True, line_spacing=1.6
    )

    # なぜ計画から始めるのか
    add_text_box(
        slide, Inches(7.0), Inches(2.0), Inches(5), Inches(0.5),
        "なぜ計画から始めるのか", font_size=20, color=ACCENT, bold=True
    )
    points = [
        "計画段階で認識のズレに気づける",
        "コードを壊すリスクがゼロ",
        "手戻りを最小限にしてから実装に入れる",
    ]
    y = Inches(2.7)
    for pt in points:
        add_shape_bg(slide, Inches(7.0), y, Inches(5.0), Inches(0.8),
                     RGBColor(0x25, 0x22, 0x20), 0.03)
        add_text_box(
            slide, Inches(7.4), y + Inches(0.15), Inches(4.2), Inches(0.5),
            pt, font_size=16, color=TEXT_LIGHT
        )
        y += Inches(0.95)

    # 操作方法
    add_text_box(
        slide, Inches(0.8), Inches(5.2), Inches(6), Inches(0.5),
        "操作方法", font_size=20, color=ACCENT, bold=True
    )
    add_multiline_text(
        slide, Inches(1.2), Inches(5.8), Inches(10), Inches(1.3),
        [
            "Shift+Tab で切替（Normal → Auto Accept → Plan）",
            "claude --permission-mode plan で起動",
        ],
        font_size=16, bullet=True, line_spacing=1.5
    )

    add_page_number(slide, 14)
    add_speaker_notes(slide,
        "Plan Modeは読み取り専用のモードです。ファイルを見たり検索はできますが、"
        "編集やコマンド実行はできません。\n"
        "なぜ計画から始めるのか。公式のベストプラクティスにも"
        "「計画ステップがないと、Claudeは直接コーディングに飛び込む傾向がある」"
        "と書かれています。\n"
        "計画なしで実装を始めると、前提の認識ズレで根本から作り直しになることがあります。\n"
        "Plan Modeなら壊すリスクはゼロです。安心して調査・計画できます。\n"
        "操作はShift+Tabで切替。起動時に --permission-mode plan で指定もできます。\n"
        "補足ですが、Extended Thinkingとは別の概念です。"
        "Extended Thinkingは深く考えさせる機能。Plan Modeはツールの制限。"
        "両方を同時に使うこともできます。")


def slide_15_plan_mode_workflow(prs):
    """Plan Mode の実践ワークフロー."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_label(slide, Inches(0.8), Inches(0.4), "PART 3 — Plan Mode")
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        "Plan Mode の実践ワークフロー", font_size=32, color=TEXT_WHITE, bold=True
    )

    steps = [
        ("Step 1", "Plan Mode で調査・計画",
         "「計画を立てて。実装はまだしないで。」",
         "Claudeがコードを読み、方針を提案する", GREEN),
        ("Step 2", "計画をレビュー",
         "「後方互換性はどうする？」「テストは？」",
         "不明点をつぶし、計画を確定する", HIGHLIGHT),
        ("Step 3", "実装を許可",
         "「OK、実装して」/ 調査だけならここで完了",
         "計画に沿って実装、または調査結果を活用", ACCENT),
        ("Step 4", "結果を確認",
         "テスト実行 → PR作成",
         "変更内容を取りまとめて完了", TEXT_WHITE),
    ]

    y = Inches(1.9)
    for step, title, instruction, desc, color in steps:
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.9),
                     RGBColor(0x25, 0x22, 0x20), 0.02)
        add_text_box(
            slide, Inches(1.2), y + Inches(0.05), Inches(1.2), Inches(0.4),
            step, font_size=14, color=color, bold=True, font_name=FONT_BODY
        )
        add_text_box(
            slide, Inches(2.5), y + Inches(0.05), Inches(3.5), Inches(0.4),
            title, font_size=17, color=TEXT_WHITE, bold=True
        )
        add_text_box(
            slide, Inches(1.2), y + Inches(0.48), Inches(5.0), Inches(0.4),
            instruction, font_size=13, color=HIGHLIGHT, font_name=FONT_CODE
        )
        add_text_box(
            slide, Inches(7.0), y + Inches(0.25), Inches(5.0), Inches(0.5),
            desc, font_size=14, color=TEXT_LIGHT
        )
        y += Inches(1.05)

    # Step 1で伝えておくと良いこと
    add_text_box(
        slide, Inches(0.8), Inches(6.0), Inches(5), Inches(0.4),
        "Step 1で伝えておくと良いこと", font_size=15, color=ACCENT, bold=True
    )
    checklist_items = [
        ("受け入れ基準", "どうなったら完了か"),
        ("禁止事項", "絶対にやってはいけない変更"),
        ("想定外の扱い", "計画にない問題を見つけたらどうするか"),
    ]
    x = Inches(0.8)
    for item, desc in checklist_items:
        add_shape_bg(slide, x, Inches(6.5), Inches(3.6), Inches(0.7),
                     RGBColor(0x25, 0x22, 0x20), 0.03)
        add_text_box(
            slide, x + Inches(0.15), Inches(6.53), Inches(3.3), Inches(0.3),
            item, font_size=13, color=GREEN, bold=True
        )
        add_text_box(
            slide, x + Inches(0.15), Inches(6.85), Inches(3.3), Inches(0.3),
            desc, font_size=11, color=TEXT_DIM
        )
        x += Inches(3.8)

    add_page_number(slide, 15)
    add_speaker_notes(slide,
        "4ステップのワークフローです。\n"
        "Step 1：「計画を立てて。実装はまだしないで。」この一言が鍵です。\n"
        "このとき、受け入れ基準や禁止事項も伝えておくとブレません。"
        "たとえば「既存APIは壊さない」「計画にない問題を見つけたら報告して止まって」など。\n"
        "Step 2：計画が返ってきたらレビューします。"
        "「後方互換性は？」「テストは？」と聞いてみてください。\n"
        "Step 3：計画に納得したら「OK、実装して」。"
        "調査だけが目的ならここで完了です。\n"
        "Step 4：テストを実行して、結果を確認してコミット。\n"
        "小さな修正でもこの流れを習慣にすると、手戻りが劇的に減ります。")


def slide_16_summary(prs):
    """まとめ / 次回予告."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
        "まとめ / 次回予告", font_size=36, color=TEXT_WHITE, bold=True
    )

    takeaways = [
        ("Claude Code はコード補完ではなく\n「作業を進める」ツール",
         "調査→計画→実装→検証→取りまとめの一連フロー"),
        ("伝え方で結果が変わる",
         "目標・背景・現状・期待出力を伝え、たたき台→改善ループで進める"),
        ("いきなりコードを書かせない",
         "Plan Modeで方針を固めてから実装に入る"),
    ]

    y = Inches(1.6)
    for i, (title, desc) in enumerate(takeaways):
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.2),
                     RGBColor(0x25, 0x22, 0x20), 0.02)
        num = str(i + 1)
        add_text_box(
            slide, Inches(1.2), y + Inches(0.1), Inches(0.5), Inches(0.5),
            num, font_size=28, color=ACCENT, bold=True, font_name=FONT_BODY
        )
        add_text_box(
            slide, Inches(2.0), y + Inches(0.1), Inches(9.5), Inches(0.5),
            title, font_size=20, color=TEXT_WHITE, bold=True
        )
        add_text_box(
            slide, Inches(2.0), y + Inches(0.65), Inches(9.5), Inches(0.5),
            desc, font_size=15, color=TEXT_LIGHT
        )
        y += Inches(1.4)

    # CTA
    add_text_box(
        slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
        "日常業務の一部を、Claude Codeに任せてみてください",
        font_size=18, color=ACCENT, bold=True
    )

    add_page_number(slide, 16)
    add_speaker_notes(slide,
        "まとめです。今日の持ち帰りは3つ。\n"
        "1つ目：Claude Codeはコード補完ではなく「作業を進める」ツール。\n"
        "2つ目：伝え方で結果が変わる。目標・背景・現状・期待出力を意識してみてください。\n"
        "3つ目：いきなりコードを書かせない。Plan Modeで方針を固めてから実装に入る。\n"
        "日常業務の一部を、まず任せてみてください。\n"
        "次回は「プロジェクトに合わせたカスタマイズ」。"
        "CLAUDE.md、Skills、Hooksを扱います。\n"
        "では、質疑応答に移ります。")


# ── Main ───────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_series_overview(prs)
    slide_03_first_step(prs)
    slide_04_what_is_cc(prs)
    slide_05_tools(prs)
    slide_06_role_comparison(prs)
    slide_07_workflow_shift(prs)
    slide_08_task_example(prs)
    slide_09_section_communication(prs)
    slide_10_agent_mental_model(prs)
    slide_11_good_instructions(prs)
    slide_12_iteration_loop(prs)
    slide_13_section_plan_mode(prs)
    slide_14_plan_mode_what(prs)
    slide_15_plan_mode_workflow(prs)
    slide_16_summary(prs)

    output_path = "slides/session1.pptx"
    prs.save(output_path)
    print(f"Generated: {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
