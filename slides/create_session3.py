# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx"]
# ///
"""Generate Claude_session3.pptx matching session1/2 design."""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens ──────────────────────────────────────────
BG_MAIN   = RGBColor(0x1B, 0x19, 0x18)
BG_SECTION = RGBColor(0x23, 0x1F, 0x1D)
CARD_BG   = RGBColor(0x2D, 0x29, 0x26)
CARD_BG2  = RGBColor(0x3A, 0x36, 0x32)
CARD_BG3  = RGBColor(0x30, 0x2D, 0x29)
CARD_BG4  = RGBColor(0x2C, 0x2A, 0x26)

ACCENT    = RGBColor(0xD9, 0x77, 0x57)  # terracotta
ACCENT2   = RGBColor(0x6A, 0x9B, 0xCC)  # blue
ACCENT3   = RGBColor(0x7B, 0xC9, 0xA0)  # green
RED       = RGBColor(0xE0, 0x70, 0x70)

WHITE     = RGBColor(0xFA, 0xF9, 0xF5)
BODY      = RGBColor(0xE0, 0xDE, 0xD8)
MUTED     = RGBColor(0xB0, 0xAE, 0xA5)
MUTED2    = RGBColor(0xC5, 0xC3, 0xBB)
DARK_MUTED = RGBColor(0x6B, 0x69, 0x63)
STRIPE_C  = RGBColor(0x2D, 0x29, 0x26)
LINE_C    = RGBColor(0x4A, 0x47, 0x42)

FONT_JP = "Hiragino Sans"
FONT_EN = "Helvetica Neue"
FONT_CODE = "SF Mono"

SLIDE_W = 12192000
SLIDE_H = 6858000
TOTAL_SLIDES = 20

# ── Helpers ────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill=None, name=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if name:
        shape.name = name
    return shape


def add_rounded_rect(slide, left, top, width, height, fill=None, name=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if name:
        shape.name = name
    shape.adjustments[0] = 0.04
    return shape


def _set_font(run, font_name, size, bold=False, color=WHITE):
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font=FONT_JP,
                size=Pt(16), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                name=None, line_spacing=None):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.fill.background()
    if name:
        txbox.name = name
    tf = txbox.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        _set_font(run, font, size, bold, color)
    return txbox


def add_part_label(slide, text):
    return add_textbox(slide, Emu(731520), Emu(365760), Emu(8000000), Emu(274320),
                       text, font=FONT_EN, size=Emu(152400), bold=True, color=ACCENT)


def add_slide_title(slide, text, y=Emu(548640)):
    return add_textbox(slide, Emu(731520), y, Emu(10000000), Emu(548640),
                       text, font=FONT_JP, size=Emu(355600), bold=True, color=WHITE)


def add_page_num(slide, num):
    return add_textbox(slide, Emu(SLIDE_W - 1200000), Emu(SLIDE_H - 400000),
                       Emu(900000), Emu(274320),
                       f"{num} / {TOTAL_SLIDES}", font=FONT_EN,
                       size=Emu(139700), color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, text, y=Emu(6200000)):
    card = add_shape(slide, Emu(731520), y, Emu(10728960), Emu(457200), fill=CARD_BG4)
    accent = add_shape(slide, Emu(731520), y, Emu(54864), Emu(457200), fill=ACCENT)
    tb = add_textbox(slide, Emu(914400), y + Emu(91440), Emu(10400000), Emu(365760),
                     text, font=FONT_JP, size=Emu(177800), bold=True, color=ACCENT)
    return card, accent, tb


def add_divider(slide, left, top, width):
    line = slide.shapes.add_connector(
        1, left, top, left + width, top
    )
    line.line.color.rgb = LINE_C
    line.line.width = Pt(0.75)
    return line


def add_code_block(slide, left, top, width, height, text):
    card = add_rounded_rect(slide, left, top, width, height, fill=RGBColor(0x1B, 0x19, 0x18))
    tb = add_textbox(slide, left + Emu(137160), top + Emu(91440),
                     width - Emu(274320), height - Emu(182880),
                     text, font=FONT_CODE, size=Emu(152400), color=BODY)
    return card, tb


def new_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def add_speaker_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def _build_table(slide, left, top, col_widths, headers, rows,
                 header_color=ACCENT, row_font=FONT_JP, code_cols=None):
    """Helper to build styled tables."""
    if code_cols is None:
        code_cols = set()
    col_x = [left]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)
    row_h = Emu(340000)
    header_h = Emu(365760)

    # Headers
    for hdr, cx, cw in zip(headers, col_x, col_widths):
        add_shape(slide, cx, top, cw, header_h, fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), top + Emu(68580), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(165100), bold=True, color=header_color)

    # Rows
    for i, row_data in enumerate(rows):
        ry = top + header_h + Emu(i * 340000)
        for j, (cell, cx, cw) in enumerate(zip(row_data, col_x, col_widths)):
            add_shape(slide, cx, ry, cw, row_h, fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j in code_cols else row_font
            fc = ACCENT if j in code_cols else BODY
            add_textbox(slide, cx + Emu(91440), ry + Emu(45720), cw - Emu(182880), Emu(274320),
                        cell, font=fn, size=Emu(152400), color=fc)
    return top + header_h + Emu(len(rows) * 340000)


# ── Slide builders ─────────────────────────────────────────

def build_title_slide(prs):
    """Slide 01: Title."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    ring = slide.shapes.add_shape(MSO_SHAPE.DONUT, Emu(8382000), Emu(-1143000),
                                   Emu(5486400), Emu(5486400))
    ring.line.fill.background()
    ring.fill.background()
    ring.line.color.rgb = RGBColor(0x3A, 0x36, 0x32)
    ring.line.width = Pt(2)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(9144000), Emu(2743200),
                                     Emu(2286000), Emu(2286000))
    circle.line.fill.background()
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT

    arc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(7620000), Emu(1600200),
                                  Emu(3200400), Emu(3200400))
    arc.line.fill.background()
    arc.fill.solid()
    arc.fill.fore_color.rgb = ACCENT

    small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(457200), Emu(5486400),
                                    Emu(914400), Emu(914400))
    small.line.fill.background()
    small.fill.background()
    small.line.color.rgb = RGBColor(0x3A, 0x36, 0x32)
    small.line.width = Pt(2)

    add_shape(slide, Emu(731520), Emu(1371600), Emu(54864), Emu(3200400), fill=ACCENT)

    for i, x in enumerate([1097280, 1371600, 1645920]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x), Emu(1097280),
                                      Emu(91440), Emu(91440))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT

    add_shape(slide, Emu(10058400), Emu(5486400), Emu(1828800), Emu(91440), fill=STRIPE_C)

    add_textbox(slide, Emu(1097280), Emu(1371600), Emu(6400800), Emu(1554480),
                "Claude Code\n大規模開発に挑む",
                font=FONT_JP, size=Emu(660400), bold=True, color=WHITE)

    add_textbox(slide, Emu(1097280), Emu(3109920), Emu(6400800), Emu(548640),
                "〜 Skills / サブエージェント / Hooks / MCP / worktree 〜",
                font=FONT_JP, size=Emu(304800), color=ACCENT)

    add_divider(slide, Emu(1097280), Emu(3840480), Emu(2743200))

    add_textbox(slide, Emu(1097280), Emu(4023360), Emu(6400800), Emu(457200),
                "社内勉強会  |  第3回 / 全3回  |  20 min",
                font=FONT_EN, size=Emu(203200), color=MUTED)

    add_speaker_notes(slide,
        "皆さんこんにちは。では第3回、最終回を始めます。\n"
        "今回は「大規模開発に挑む」がテーマです。\n\n"
        "前回まで、第1回ではClaude Codeの基本的な使い方とコミュニケーションの考え方、"
        "第2回ではCLAUDE.mdやsettings.jsonによるプロジェクトセットアップを学びました。\n"
        "今回は、そのセットアップされた環境をベースに、より大きな開発に対応するための"
        "実践的な使い方に踏み込んでいきます。")


def build_roadmap_slide(prs):
    """Slide 02: Roadmap (3 cards)."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_textbox(slide, Emu(736600), Emu(457200), Emu(7620000), Emu(635000),
                "全3回のロードマップ", font=FONT_JP, size=Emu(406400), bold=True, color=WHITE)

    card_w = Emu(3429000)
    card_h = Emu(3937000)
    card_y = Emu(1270000)
    gap = Emu(228600)

    card_data = [
        {"num": "1", "label": "第1回（済）", "title": "理解する",
         "desc": "基本 / コミュニケーション / Plan Mode",
         "accent": MUTED, "card_fill": CARD_BG3, "label_color": MUTED},
        {"num": "2", "label": "第2回（済）", "title": "プロジェクトに\n定着させる",
         "desc": "CLAUDE.md / settings.json",
         "accent": MUTED, "card_fill": CARD_BG3, "label_color": MUTED},
        {"num": "3", "label": "第3回（今日）", "title": "大規模開発に挑む",
         "desc": "Skills / サブエージェント /\nHooks / MCP / worktree",
         "accent": ACCENT, "card_fill": CARD_BG2, "label_color": ACCENT},
    ]

    for i, d in enumerate(card_data):
        x = Emu(736600) + i * (card_w + gap)

        add_rounded_rect(slide, x, card_y, card_w, card_h, fill=d["card_fill"])
        add_shape(slide, x, card_y, card_w, Emu(63500), fill=d["accent"])

        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Emu(228600), card_y + Emu(304800),
                                            Emu(609600), Emu(609600))
        num_shape.line.fill.background()
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = d["accent"]
        tf = num_shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = d["num"]
        _set_font(run, FONT_EN, Emu(254000), True, WHITE)

        add_textbox(slide, x + Emu(965200), card_y + Emu(355600), Emu(2286000), Emu(508000),
                    d["label"], font=FONT_JP, size=Emu(203200), bold=True, color=d["label_color"])

        add_textbox(slide, x + Emu(228600), card_y + Emu(1143000), Emu(2984500), Emu(1143000),
                    d["title"], font=FONT_JP, size=Emu(279400), bold=True, color=WHITE)

        add_divider(slide, x + Emu(228600), card_y + Emu(2476500), Emu(2984500))

        add_textbox(slide, x + Emu(228600), card_y + Emu(2654300), Emu(2984500), Emu(1016000),
                    d["desc"], font=FONT_JP, size=Emu(190500), color=MUTED2)

    add_textbox(slide, Emu(736600), Emu(5500000), Emu(10728960), Emu(457200),
                "今日のゴール: 繰り返しの指示を仕組み化し、複数タスクの並列開発ができるようになる",
                font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    add_page_num(slide, 2)

    add_speaker_notes(slide,
        "全3回の最終回です。\n\n"
        "第1回では基本とコミュニケーション、第2回ではCLAUDE.mdとsettings.jsonによるセットアップを扱いました。\n\n"
        "今回は、大規模開発に対応するための5つのトピックを扱います。\n"
        "この順番には意味がありまして、まずSkillsとサブエージェントで指示を仕組み化し、"
        "次にHooks・MCPで自動化と外部連携を整えて、最後にworktreeでそれらすべてを束ねて"
        "並列開発を実現する、という積み上げ式の流れでいきます。")


def build_agenda_slide(prs):
    """Slide 03: Agenda (3 cards)."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_textbox(slide, Emu(736600), Emu(457200), Emu(7620000), Emu(635000),
                "今日のアジェンダ", font=FONT_JP, size=Emu(406400), bold=True, color=WHITE)

    card_w = Emu(3429000)
    card_h = Emu(3600000)
    card_y = Emu(1270000)
    gap = Emu(228600)

    cards = [
        {"part": "PART 1", "title": "Skills &\nサブエージェント",
         "desc": "Claudeの能力を拡張する2つの仕組み。\n繰り返しの指示を仕組み化し、\n重い作業を委譲する",
         "accent": ACCENT},
        {"part": "PART 2", "title": "Hooks & MCP",
         "desc": "品質チェックの自動化と\n外部サービス連携。\n短めにいきます",
         "accent": ACCENT2},
        {"part": "PART 3", "title": "worktree",
         "desc": "PART 1・2で整えた環境を束ねて、\n複数タスクを並列で進める。\nデモあり",
         "accent": ACCENT3},
    ]

    for i, d in enumerate(cards):
        x = Emu(736600) + i * (card_w + gap)
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill=CARD_BG)
        add_shape(slide, x, card_y, card_w, Emu(63500), fill=d["accent"])

        add_textbox(slide, x + Emu(228600), card_y + Emu(228600), Emu(2984500), Emu(365760),
                    d["part"], font=FONT_EN, size=Emu(203200), bold=True, color=d["accent"])

        add_textbox(slide, x + Emu(228600), card_y + Emu(700000), Emu(2984500), Emu(900000),
                    d["title"], font=FONT_JP, size=Emu(330200), bold=True, color=WHITE)

        add_divider(slide, x + Emu(228600), card_y + Emu(1800000), Emu(2984500))

        add_textbox(slide, x + Emu(228600), card_y + Emu(1950000), Emu(2984500), Emu(1400000),
                    d["desc"], font=FONT_JP, size=Emu(177800), color=MUTED2)

    add_page_num(slide, 3)

    add_speaker_notes(slide,
        "今日のアジェンダです。\n\n"
        "PART 1はSkillsとサブエージェント。ここが今日のメインで、厚めにやります。"
        "繰り返しの指示を仕組み化する方法と、重い作業を委譲する方法の話です。\n\n"
        "PART 2はHooksとMCP。品質チェックの自動化と外部連携の話で、こちらは手短にいきます。\n\n"
        "PART 3はworktree。PART 1・2で整えた環境をすべて束ねて、並列開発を実現するという話です。"
        "最後にデモもお見せしますので、お楽しみにしてください。")


def build_section_slide(prs, part_num, title, subtitle, slide_num):
    """Section divider slide."""
    slide = new_slide(prs)
    set_bg(slide, BG_SECTION)

    line_y = Emu(SLIDE_H // 2)
    add_divider(slide, Emu(731520), line_y, Emu(1828800))
    add_divider(slide, Emu(SLIDE_W - 731520 - 1828800), line_y, Emu(1828800))

    add_textbox(slide, Emu(0), Emu(2000000), Emu(SLIDE_W), Emu(457200),
                f"PART {part_num}", font=FONT_EN, size=Emu(228600), bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(slide, Emu(0), Emu(2600000), Emu(SLIDE_W), Emu(762000),
                title, font=FONT_JP, size=Emu(508000), bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Emu(0), Emu(3500000), Emu(SLIDE_W), Emu(457200),
                subtitle, font=FONT_JP, size=Emu(304800),
                color=MUTED2, align=PP_ALIGN.CENTER)

    add_page_num(slide, slide_num)

    notes_map = {
        1: ("では早速PART 1に入りましょう。Skillsとサブエージェントです。\n\n"
            "この2つは結構つながりが深くて、Skillsで指示を仕組み化し、サブエージェントで重い作業を委譲する。"
            "どちらもClaudeの能力を拡張する仕組みなんですが、使いどころが違います。"
            "そのあたりを順番に見ていきます。"),
        2: ("ここからPART 2です。HooksとMCPを扱います。\n\n"
            "PART 1でSkillsとサブエージェントという「Claudeの能力を拡張する仕組み」を見ました。"
            "PART 2ではもう少し仕組みの側の話で、Hooksで品質チェックを自動化し、"
            "MCPで外部サービスと連携する方法を紹介します。どちらも手短にいきますね。"),
        3: ("PART 3、worktreeです。ここが今日の話の集大成になります。\n\n"
            "PART 1でSkillsとサブエージェント、PART 2でHooksとMCPを見てきました。"
            "これらの仕組みをすべて使いこなした上で、複数のタスクを並列で進める。"
            "それがworktreeとの組み合わせです。\nでは具体的にどうやるのか見ていきましょう。"),
    }
    if part_num in notes_map:
        add_speaker_notes(slide, notes_map[part_num])


def build_skill_intro_slide(prs):
    """Slide 05: What are Skills."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "スキルとは")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "スラッシュコマンドとして使える拡張機能",
                font=FONT_JP, size=Emu(228600), bold=True, color=ACCENT)

    # Code tree
    add_code_block(slide, Emu(731520), Emu(1700000), Emu(4800000), Emu(900000),
                   ".claude/skills/\n"
                   "├── review/\n"
                   "│   └── SKILL.md      # /review で呼び出せる\n"
                   "└── deploy.md          # /deploy で呼び出せる")

    add_textbox(slide, Emu(731520), Emu(2750000), Emu(5000000), Emu(365760),
                "自作スキルも .claude/skills/ に置くだけで追加できる",
                font=FONT_JP, size=Emu(177800), color=MUTED)

    # Comparison table on right
    add_textbox(slide, Emu(6000000), Emu(1200000), Emu(5000000), Emu(365760),
                "CLAUDE.mdとの違い", font=FONT_JP, size=Emu(228600), bold=True, color=WHITE)

    col_w = [Emu(1400000), Emu(2200000), Emu(2200000)]
    col_x = [Emu(6000000), Emu(7400000), Emu(9600000)]
    headers = ["", "CLAUDE.md", "スキル"]
    hy = Emu(1650000)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(340000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), hy + Emu(60000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(152400), bold=True, color=ACCENT)

    rows = [
        ("読み込み", "常に自動", "必要なときだけ"),
        ("内容", "基本ルール", "特定の作業手順"),
        ("目安", "500行以下", "必要に応じて分割"),
        ("例", "コーディング規約", "レビュー手順"),
    ]
    for i, (label, c1, c2) in enumerate(rows):
        ry = hy + Emu(340000) + Emu(i * 320000)
        for j, (cell, cx, cw) in enumerate(zip([label, c1, c2], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(310000), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fc = WHITE if j == 0 else BODY
            fb = j == 0
            add_textbox(slide, cx + Emu(68580), ry + Emu(50000), cw - Emu(137160), Emu(274320),
                        cell, font=FONT_JP, size=Emu(139700), bold=fb, color=fc)

    # Footer
    add_footer(slide, "→ CLAUDE.md = 常識、スキル = 専門知識。Agent Skills Open Standardに準拠")
    add_page_num(slide, 5)

    add_speaker_notes(slide,
        "まずスキルの話から入るんですが、その前に前回の内容を少し振り返ります。\n"
        "前回、CLAUDE.mdが大きくなりすぎないように注意しましょう、という話をしました。"
        "500行以下が目安でしたよね。\n"
        "じゃあ、CLAUDE.mdに書ききれない専門知識や作業手順はどうするのか。そこでスキルの出番です。\n\n"
        "スキルは、スラッシュコマンドとして呼び出せる拡張機能です。\n"
        "最初から使えるコマンドやスキルも多数用意されています。これらは後ほどまとめて紹介します。"
        "今日は主にカスタムスキルの作り方を紹介しますが、最後に組み込みの /batch がworktreeと"
        "組み合わさるとどうなるか、という話もします。\n\n"
        "前回学んだCLAUDE.mdとの違いが大事です。CLAUDE.mdは「常に」自動で読み込まれる「常識」。"
        "スキルは「必要なときだけ」読み込まれる「専門知識」です。\n\n"
        "CLAUDE.mdに全部詰め込むとコンテキストを圧迫するので、状況に応じて読み込ませたい情報は"
        "スキルに分けるのが良いかなと思います。\n\n"
        "ちなみに、スキルはAgent Skills Open Standardという業界標準に準拠しています。"
        "Claude Code用に作ったスキルが、GitHub CopilotやOpenAI Codex CLIでもそのまま動きます。"
        "学習コストが一度で済むのも嬉しいポイントです。")


def build_skill_types_slide(prs):
    """Slide 06: Two types of skills."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "スキルの2タイプ")

    half_w = Emu(5200000)
    lx = Emu(731520)
    rx = Emu(6200000)

    # Left: Reference content
    add_textbox(slide, lx, Emu(1200000), half_w, Emu(365760),
                "参照コンテンツ — 知識を渡す", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    add_code_block(slide, lx, Emu(1600000), Emu(5000000), Emu(2400000),
                   "---\nname: api-conventions\n"
                   "description: APIエンドポイント作成時の\n"
                   "             設計パターン\n---\n\n"
                   "APIエンドポイントを書くときは以下に従う。\n"
                   "- RESTfulな命名規則を使う\n"
                   "- 一貫したエラーフォーマットを返す\n"
                   "- リクエストバリデーションを含める")

    add_textbox(slide, lx, Emu(4100000), half_w, Emu(365760),
                "→ 関連する作業を検知したとき、自動で読み込まれる",
                font=FONT_JP, size=Emu(165100), bold=True, color=ACCENT)

    # Right: Task content
    add_textbox(slide, rx, Emu(1200000), half_w, Emu(365760),
                "タスクコンテンツ — 手順を渡す", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT2)

    add_code_block(slide, rx, Emu(1600000), Emu(5000000), Emu(2400000),
                   "---\nname: issue-planner\n"
                   "description: Issueの調査と設計を行う\n"
                   "disable-model-invocation: true\n---\n\n"
                   "Issue #$ARGUMENTS の設計を行う。\n"
                   "1. gh issue viewでIssue詳細を取得\n"
                   "2. 関連コードを検索し影響範囲を特定\n"
                   "3. 実装方針をまとめる\n"
                   "4. ISSUE.mdに設計メモを追記")

    add_textbox(slide, rx, Emu(4100000), half_w, Emu(365760),
                "→ /issue-planner 42 のように明示的に呼び出す",
                font=FONT_JP, size=Emu(165100), bold=True, color=ACCENT2)

    add_page_num(slide, 6)

    add_speaker_notes(slide,
        "スキルには大きく2つのタイプがあります。\n\n"
        "参照コンテンツは「知識」を渡すタイプです。API設計パターンのように、"
        "関連する作業をClaudeが検知したときに自動で読み込まれます。\n\n"
        "タスクコンテンツは「手順」を渡すタイプ。/issue-planner 42 のように明示的に呼び出して使います。\n"
        "disable-model-invocation: true を設定すると、自動では呼び出されず、"
        "スラッシュコマンドとして呼び出したときだけ読み込まれるようになります。\n\n"
        "$ARGUMENTS というプレースホルダーで引数を受け取れるので、たとえばIssue番号を渡すと"
        "調査から設計まで一気に進みます。\n"
        "さらに !`command` という構文を使うと、シェルコマンドの実行結果をスキル本文に動的に注入できます。"
        "たとえば !`git branch --show-current` と書けば、現在のブランチ名がスキル実行時に展開されます。\n\n"
        "こういった機能を組み合わせると、たとえばissue-workerでworktreeを作った後、"
        "issue-plannerで設計を固めてから実装に入る、といった連携もできるようになります。")


def build_skill_examples_slide(prs):
    """Slide 07: Practical skill examples."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "実践的なスキル例")

    half_w = Emu(5200000)
    lx = Emu(731520)
    rx = Emu(6200000)

    # Left: issue-worker
    add_textbox(slide, lx, Emu(1200000), half_w, Emu(365760),
                "例1: Issue作業スキル", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    add_code_block(slide, lx, Emu(1600000), Emu(5000000), Emu(2200000),
                   "---\nname: issue-worker\n"
                   "description: GitHub Issueを分析し、\n"
                   "  worktreeで作業開始まで導く\n"
                   "disable-model-invocation: true\n---\n\n"
                   "$ARGUMENTS のIssue番号について以下を実行。\n"
                   "1. gh issue view で Issue 詳細を取得\n"
                   "2. git worktree add で作業ブランチ作成\n"
                   "3. Issue内容から ISSUE.md を自動生成\n"
                   "4. 新しいworktreeでClaude Code起動")

    add_textbox(slide, lx, Emu(3900000), half_w, Emu(274320),
                "使い方: /issue-worker 42",
                font=FONT_CODE, size=Emu(165100), color=ACCENT)

    # Right: pr-creator
    add_textbox(slide, rx, Emu(1200000), half_w, Emu(365760),
                "例2: PR作成スキル", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT2)

    add_code_block(slide, rx, Emu(1600000), Emu(5000000), Emu(2200000),
                   "---\nname: pr-creator\n"
                   "description: チームのテンプレートに\n"
                   "  沿ってPRを作成する\n"
                   "disable-model-invocation: true\n---\n\n"
                   "PRを作成する。\n"
                   "1. git diff で変更内容を確認\n"
                   "2. テンプレートに沿って本文を作成\n"
                   "   - ## 概要 / ## 背景 / ## 影響範囲\n"
                   "3. gh pr create で提出")

    add_textbox(slide, rx, Emu(3900000), half_w, Emu(274320),
                "使い方: /pr-creator",
                font=FONT_CODE, size=Emu(165100), color=ACCENT2)

    # Footer
    add_footer(slide,
               "→ 同じ指示を何度も書いている → スキル化を検討。.claude/skills/ をgit管理でチーム共有",
               y=Emu(4400000))
    add_page_num(slide, 7)

    add_speaker_notes(slide,
        "チームで共有すると嬉しいスキル例を2つ紹介します。\n\n"
        "1つ目は issue-worker。Issue番号を渡すだけで、worktreeの作成からISSUE.mdの生成まで自動で行います。"
        "この後のPART 3で紹介するworktreeとの組み合わせで、複数Issueの並列開発ができるようになるので、"
        "ここは後でまた出てきます。\n\n"
        "2つ目はPR作成。PRのフォーマットがバラバラだと、レビューする側も大変ですよね。"
        "テンプレートをスキル化しておけば、チーム全体のPR品質が揃います。\n\n"
        "どちらも .claude/skills/ に置いてgit管理すれば、チーム全員が同じスキルを使えます。"
        "スキルを作るタイミングとしては、同じ指示を何度も手打ちしていると感じたとき。"
        "逆に一度きりの作業なら、急いでスキル化する必要はありません。")


def build_skill_design_slide(prs):
    """Slide 08: Skill design points."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "スキルの設計ポイント")

    # Settings fields table
    add_textbox(slide, Emu(731520), Emu(1100000), Emu(5000000), Emu(365760),
                "設定フィールド", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    col_w = [Emu(2800000), Emu(800000), Emu(3600000)]
    col_x = [Emu(731520), Emu(3531520), Emu(4331520)]
    headers = ["フィールド", "必須", "説明"]
    hy = Emu(1450000)
    rh = Emu(290000)

    for hdr, cx, cw in zip(headers, col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(320000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), hy + Emu(55000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT)

    fields = [
        ("name", "Yes", "一意の識別子。/name で呼び出せる"),
        ("description", "Yes", "自動読み込み判断に使用"),
        ("allowed-tools", "No", "使用可能なツールを制限"),
        ("disable-model-invocation", "No", "trueで自動呼び出し無効"),
        ("argument-hint", "No", "引数のヒントを補完表示"),
        ("model", "No", "実行時のモデルを指定"),
        ("context", "No", "forkでサブエージェント実行"),
    ]
    for i, (field, req, desc) in enumerate(fields):
        ry = hy + Emu(320000) + Emu(i * rh)
        for j, (cell, cx, cw) in enumerate(zip([field, req, desc], col_x, col_w)):
            add_shape(slide, cx, ry, cw, rh, fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j == 0 else FONT_JP
            fc = ACCENT if j == 0 else BODY
            add_textbox(slide, cx + Emu(68580), ry + Emu(40000), cw - Emu(137160), Emu(274320),
                        cell, font=fn, size=Emu(127000), color=fc)

    # Storage table (right side)
    add_textbox(slide, Emu(8200000), Emu(1100000), Emu(3500000), Emu(365760),
                "保存場所", font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT2)

    scol_w = [Emu(1800000), Emu(1900000)]
    scol_x = [Emu(8200000), Emu(10000000)]
    for hdr, cx, cw in zip(["場所", "スコープ"], scol_x, scol_w):
        add_shape(slide, cx, Emu(1450000), cw, Emu(320000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), Emu(1505000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT2)

    storage = [
        (".claude/skills/", "プロジェクト"),
        ("~/.claude/skills/", "全プロジェクト"),
    ]
    for i, (loc, scope) in enumerate(storage):
        ry = Emu(1770000) + Emu(i * 310000)
        for j, (cell, cx, cw) in enumerate(zip([loc, scope], scol_x, scol_w)):
            add_shape(slide, cx, ry, cw, Emu(300000), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j == 0 else FONT_JP
            fc = ACCENT if j == 0 else BODY
            add_textbox(slide, cx + Emu(68580), ry + Emu(45000), cw - Emu(137160), Emu(274320),
                        cell, font=fn, size=Emu(127000), color=fc)

    # Design tips
    add_divider(slide, Emu(731520), Emu(3800000), Emu(10728960))

    add_textbox(slide, Emu(731520), Emu(3950000), Emu(5000000), Emu(365760),
                "設計のコツ", font=FONT_JP, size=Emu(203200), bold=True, color=WHITE)

    tips = [
        "descriptionが鍵。ここを読んでClaudeが自動読み込みを判断する",
        "スキルの中身はClaudeへの指示そのもの。人に頼むように書く",
        "1スキル1責務。大きくなったら分割する",
    ]
    ty = Emu(4350000)
    for tip in tips:
        add_textbox(slide, Emu(914400), ty, Emu(10000000), Emu(274320),
                    f"•  {tip}", font=FONT_JP, size=Emu(177800), color=BODY)
        ty += Emu(340000)

    add_footer(slide, "→ ここまでが「スキル = メイン会話の中で動く仕組み」。次は「独立して動く」サブエージェントへ",
               y=Emu(5500000))
    add_page_num(slide, 8)

    add_speaker_notes(slide,
        "スキルの設計で押さえておくべきポイントをいくつかお伝えします。\n\n"
        "まず設定フィールドですが、全部覚える必要はなくて、よく使うものだけ紹介します。\n\n"
        "allowed-tools で使えるツールを制限できます。たとえばレビュー用スキルなら読み取り系ツールだけに"
        "絞る、といった使い方です。argument-hint を設定すると、スラッシュコマンド入力時に引数のヒントが"
        "表示されて使いやすくなります。\n\n"
        "context: fork を設定すると、スキルが独立したサブエージェントのコンテキストで実行されます。"
        "重い処理でもメインの会話を圧迫しません。これは次に紹介するサブエージェントの特性をスキルに"
        "持たせる設定、という感じです。\n\n"
        "ちなみに user-invocable: false という設定もあって、これを使うとメニューには表示されず、"
        "Claudeが関連する作業を検知したときだけ自動で読み込まれる背景知識として機能します。\n\n"
        "保存場所はプロジェクトとユーザーの2つ。チーム共有ならプロジェクト側、個人用ならユーザー側に置きます。"
        "前回のCLAUDE.mdと同じ考え方ですね。\n\n"
        "ひとつ実用的な情報をお伝えしますと、skill-creator というスキル作成を支援するツールも用意されています。"
        "対話形式でスキルを作成し、テストやベンチマークで品質を検証できます。スキルが増えてきたら活用してみてください。\n\n"
        "ここまでがスキルの話です。スキルは「メイン会話の中で動く仕組み」。"
        "次に紹介するサブエージェントは「独立して動く仕組み」です。この違いが大事になってきます。")


def build_builtin_commands_slide(prs):
    """Slide 09: Built-in commands & skills."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "最初から使えるコマンド & スキル")

    # Top table: builtin commands
    add_textbox(slide, Emu(731520), Emu(1100000), Emu(10000000), Emu(365760),
                "ビルトインコマンド（固定ロジック）— 即座に結果を返す",
                font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT)

    col_w = [Emu(2800000), Emu(7928960)]
    col_x = [Emu(731520), Emu(3531520)]
    hy = Emu(1450000)
    rh = Emu(280000)

    for hdr, cx, cw in zip(["コマンド", "用途"], col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(310000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), hy + Emu(55000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT)

    cmds = [
        ("/security-review", "ブランチの変更をセキュリティ脆弱性分析"),
        ("/pr-comments", "GitHub PRのコメントを取得・表示"),
        ("/diff", "差分のインタラクティブビューア"),
        ("/compact", "コンテキストを圧縮して空きを作る"),
        ("/init", "CLAUDE.md を対話的に初期化"),
        ("他50以上", "/help で一覧表示"),
    ]
    for i, (cmd, usage) in enumerate(cmds):
        ry = hy + Emu(310000) + Emu(i * rh)
        for j, (cell, cx, cw) in enumerate(zip([cmd, usage], col_x, col_w)):
            add_shape(slide, cx, ry, cw, rh, fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j == 0 else FONT_JP
            fc = ACCENT if j == 0 else BODY
            add_textbox(slide, cx + Emu(68580), ry + Emu(40000), cw - Emu(137160), Emu(274320),
                        cell, font=fn, size=Emu(127000), color=fc)

    # Bottom table: bundled skills
    skill_y = Emu(3450000)
    add_textbox(slide, Emu(731520), skill_y, Emu(10000000), Emu(365760),
                "バンドルスキル（プロンプトベース）— エージェントを起動し複雑なタスクを遂行",
                font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT2)

    shy = skill_y + Emu(380000)
    for hdr, cx, cw in zip(["スキル", "用途"], col_x, col_w):
        add_shape(slide, cx, shy, cw, Emu(310000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), shy + Emu(55000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT2)

    skills = [
        ("/batch", "大規模変更を5〜30ユニットに分解、worktreeで並列実行→PR作成"),
        ("/simplify", "3つのレビューエージェントが品質・効率・再利用性を並列検証"),
        ("/loop", "定期的にプロンプト実行（cron対応、最大3日間）"),
        ("/debug", "セッションのデバッグログを分析"),
        ("/claude-api", "Claude API / Agent SDKリファレンスを自動読み込み"),
    ]
    for i, (sk, usage) in enumerate(skills):
        ry = shy + Emu(310000) + Emu(i * rh)
        for j, (cell, cx, cw) in enumerate(zip([sk, usage], col_x, col_w)):
            add_shape(slide, cx, ry, cw, rh, fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j == 0 else FONT_JP
            fc = ACCENT2 if j == 0 else BODY
            add_textbox(slide, cx + Emu(68580), ry + Emu(40000), cw - Emu(137160), Emu(274320),
                        cell, font=fn, size=Emu(127000), color=fc)

    add_footer(slide, "→ 特に /batch と /simplify が強力。PART 3のworktreeで再登場する")
    add_page_num(slide, 9)

    add_speaker_notes(slide,
        "ここまでカスタムスキルの作り方を見てきましたが、最初から使えるものも多数あります。\n\n"
        "/ を打つと一覧が出ますが、大きく2種類に分かれています。\n\n"
        "ビルトインコマンドは、Claude Code本体に組み込まれた固定機能です。"
        "/security-review でセキュリティチェック、/diff で差分確認、/compact でコンテキスト圧縮。"
        "50以上あるので /help で確認してみてください。これらはClaude Code本体の機能なので、"
        "ユーザーが新たに作ることはできません。\n\n"
        "「じゃあ自作コマンドは作れないの？」と思うかもしれませんが、かつてあったカスタムコマンド"
        "（.claude/commands/）は、先ほど紹介したスキルに統合されました。"
        "つまり、自分で拡張するならスキルを使う、というのが今の形です。\n\n"
        "バンドルスキルはプロンプトベースで、エージェントを起動して複雑なタスクを遂行します。"
        "特に /batch と /simplify が強力です。/batch は大規模変更をworktreeで並列実行、"
        "/simplify は3つのレビューエージェントが品質を検証します。"
        "この2つはPART 3のworktreeの話で改めて触れます。")


def build_subagent_intro_slide(prs):
    """Slide 10: What are sub-agents."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "サブエージェントとは")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "独立したコンテキストで動く専門エージェント",
                font=FONT_JP, size=Emu(228600), bold=True, color=ACCENT)

    # Flow diagram
    main_y = Emu(1800000)
    add_rounded_rect(slide, Emu(731520), main_y, Emu(4400000), Emu(457200), fill=CARD_BG2)
    add_textbox(slide, Emu(831520), main_y + Emu(91440), Emu(4200000), Emu(274320),
                "メインエージェント", font=FONT_JP, size=Emu(203200), bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, Emu(2500000), main_y + Emu(457200), Emu(800000), Emu(365760),
                "指示 ▼", font=FONT_JP, size=Emu(152400), color=ACCENT, align=PP_ALIGN.CENTER)

    sub_y = main_y + Emu(914400)
    add_rounded_rect(slide, Emu(731520), sub_y, Emu(4400000), Emu(800000), fill=CARD_BG)
    add_shape(slide, Emu(731520), sub_y, Emu(54864), Emu(800000), fill=ACCENT2)
    add_textbox(slide, Emu(914400), sub_y + Emu(68580), Emu(4000000), Emu(274320),
                "サブエージェント（独立コンテキスト）",
                font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT2)
    add_textbox(slide, Emu(914400), sub_y + Emu(365760), Emu(4000000), Emu(365760),
                "数十ファイルを検索・分析...",
                font=FONT_JP, size=Emu(177800), color=BODY)

    add_textbox(slide, Emu(2500000), sub_y + Emu(800000), Emu(800000), Emu(365760),
                "▲ 要約だけ返る", font=FONT_JP, size=Emu(152400), color=ACCENT3, align=PP_ALIGN.CENTER)

    result_y = sub_y + Emu(1280000)
    add_rounded_rect(slide, Emu(731520), result_y, Emu(4400000), Emu(457200), fill=CARD_BG2)
    add_textbox(slide, Emu(831520), result_y + Emu(91440), Emu(4200000), Emu(274320),
                "メインのコンテキストを圧迫しない",
                font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT3,
                align=PP_ALIGN.CENTER)

    # Right side: key points
    add_textbox(slide, Emu(5600000), Emu(1800000), Emu(6000000), Emu(365760),
                '本質: コンテキストの分離', font=FONT_JP,
                size=Emu(254000), bold=True, color=WHITE)

    key_points = [
        "サブエージェントに委譲したタスクは\n独立して処理される",
        "作業結果の要約だけが\nメイン会話に返される",
        "メイン会話のコンテキストを消費しない",
    ]
    ky = Emu(2300000)
    for kp in key_points:
        add_rounded_rect(slide, Emu(5600000), ky, Emu(5860960), Emu(548640), fill=CARD_BG)
        add_textbox(slide, Emu(5800000), ky + Emu(68580), Emu(5460960), Emu(457200),
                    kp, font=FONT_JP, size=Emu(177800), color=BODY)
        ky += Emu(640000)

    add_footer(slide, "→ 大量のファイル検索や広範な調査をしてもメインの作業に影響しない")
    add_page_num(slide, 10)

    add_speaker_notes(slide,
        "ここを理解しておくと、このあとのworktreeの話がすんなり入ってきますので、少しだけ丁寧にいきます。\n\n"
        "サブエージェントは、独立したコンテキストで動く専門エージェントです。\n\n"
        "たとえば「セキュリティの問題がないか調べて」と指示すると、サブエージェントが数十のファイルを"
        "読んで調査しますが、メインの会話に返されるのは要約だけです。\n\n"
        "前回学んだように、コンテキストは有限な資源でしたよね。大量のファイルを検索する作業をメインでやると"
        "すぐにコンテキストが埋まりますが、サブエージェントに委譲すればメインのコンテキストを守れます。"
        "いわば、調査を任せて報告だけもらう、という使い方です。")


def build_subagent_combined_slide(prs):
    """Slide 11: Built-in & custom sub-agents."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "組み込み & カスタムサブエージェント")

    # Mini table: 3 built-in agents
    add_textbox(slide, Emu(731520), Emu(1100000), Emu(5000000), Emu(365760),
                "3つの組み込みサブエージェント",
                font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    col_w = [Emu(2000000), Emu(1600000), Emu(3600000)]
    col_x = [Emu(731520), Emu(2731520), Emu(4331520)]
    hy = Emu(1450000)
    rh = Emu(300000)

    for hdr, cx, cw in zip(["名前", "モデル", "用途"], col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(310000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), hy + Emu(55000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT)

    agents = [
        ("Explore", "Haiku（高速）", "コードベースの検索・分析（読み取り専用）"),
        ("Plan", "—", "Planモード中の調査"),
        ("general-purpose", "Sonnet", "複雑なマルチステップタスク（全ツール利用可）"),
    ]
    for i, (name, model, usage) in enumerate(agents):
        ry = hy + Emu(310000) + Emu(i * rh)
        for j, (cell, cx, cw) in enumerate(zip([name, model, usage], col_x, col_w)):
            add_shape(slide, cx, ry, cw, rh, fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_EN if j == 0 else FONT_JP
            fc = ACCENT if j == 0 else BODY
            add_textbox(slide, cx + Emu(68580), ry + Emu(40000), cw - Emu(137160), Emu(274320),
                        cell, font=fn, size=Emu(127000), color=fc)

    add_textbox(slide, Emu(731520), Emu(2500000), Emu(7000000), Emu(274320),
                "→ Claudeが状況に応じて自動的に使い分ける",
                font=FONT_JP, size=Emu(165100), bold=True, color=MUTED)

    # Custom agent code block (right area / below)
    add_textbox(slide, Emu(731520), Emu(2900000), Emu(5000000), Emu(365760),
                "カスタムサブエージェント — .claude/agents/ にMarkdownを配置",
                font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT2)

    add_code_block(slide, Emu(731520), Emu(3300000), Emu(6000000), Emu(2600000),
                   "---\nname: code-reviewer\n"
                   "description: コードレビューの専門家。\n"
                   "             コード変更後に積極的にレビュー。\n"
                   "tools: Read, Grep, Glob, Bash\n"
                   "model: sonnet\n"
                   "memory: project\n---\n\n"
                   "あなたはシニアコードレビューアーです。\n\n"
                   "呼び出されたときは以下を実行する。\n"
                   "1. git diffで最近の変更を確認\n"
                   "2. 変更されたファイルに焦点を当てる\n"
                   "3. 即座にレビューを開始")

    # Tips on right
    tips = [
        "descriptionが重要。Claudeはこの説明を\n見て委譲先を判断する",
        "isolation: worktree を追加すると\n独立したworktreeで動作する",
        "memory: project でセッションをまたいで\n知見を蓄積してくれる",
    ]
    ty = Emu(3300000)
    for tip in tips:
        add_rounded_rect(slide, Emu(7000000), ty, Emu(4460960), Emu(500000), fill=CARD_BG)
        add_textbox(slide, Emu(7150000), ty + Emu(50000), Emu(4160960), Emu(420000),
                    tip, font=FONT_JP, size=Emu(152400), color=BODY)
        ty += Emu(560000)

    add_page_num(slide, 11)

    add_speaker_notes(slide,
        "組み込みのサブエージェントが3つあります。Exploreは高速な検索特化型、Planはプランモード中の調査用、"
        "general-purposeは汎用タイプです。これらはClaudeが状況に応じて自動的に使い分けてくれるので、"
        "皆さんが意識する必要はありません。\n\n"
        "で、カスタムのサブエージェントも作れます。.claude/agents/ にMarkdownファイルを置くだけです。"
        "スキルと同じような配置方法ですね。\n\n"
        "ここで大事なのがdescriptionです。たとえば「コード変更後に積極的にレビュー」と書いておけば、"
        "コード変更があったときに自動的にこのサブエージェントが呼ばれるようになります。\n\n"
        "ちなみに、フロントマターに isolation: worktree を追加すると、サブエージェントが独立した"
        "worktreeで動きます。ファイルの競合を気にせず並列で動かせるので、大規模な変更に有効です。"
        "この詳細はPART 3でまた触れます。\n\n"
        "もう1つ面白い設定として memory: project があります。これを設定すると、サブエージェントが"
        "セッションをまたいで知見を蓄積してくれます。レビューで見つけたパターンや頻出する問題を覚えて"
        "くれるので、使うほど賢くなるという感じです。スコープは project（チーム共有）、"
        "user（個人・全プロジェクト共通）、local（個人・プロジェクト固有）の3つがあります。")


def build_skill_vs_subagent_slide(prs):
    """Slide 12: Skill vs Sub-agent."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 1 — Skills & サブエージェント")
    add_slide_title(slide, "スキル vs サブエージェント")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "どちらもClaudeの能力を拡張するが、動作が異なる",
                font=FONT_JP, size=Emu(203200), color=MUTED)

    # Comparison table
    col_w = [Emu(2600000), Emu(4064480), Emu(4064480)]
    col_x = [Emu(731520)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    headers = ["", "スキル", "サブエージェント"]
    hy = Emu(1700000)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(365760), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), hy + Emu(68580), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT)

    rows = [
        ("コンテキスト", "メイン会話内", "独立（要約だけ返る）"),
        ("向いている作業", "知識参照、軽量なワークフロー", "大量の探索、並列実行"),
        ("配置場所", ".claude/skills/", ".claude/agents/"),
        ("設定ファイル", "SKILL.md", "Markdown"),
    ]
    for i, (label, skill, agent) in enumerate(rows):
        ry = hy + Emu(365760) + Emu(i * 370000)
        for j, (cell, cx, cw) in enumerate(zip([label, skill, agent], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(365760), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j > 0 and i >= 2 else FONT_JP
            fc = WHITE if j == 0 else BODY
            fb = j == 0
            add_textbox(slide, cx + Emu(91440), ry + Emu(68580), cw - Emu(182880), Emu(274320),
                        cell, font=fn, size=Emu(177800), bold=fb, color=fc)

    add_divider(slide, Emu(731520), Emu(3600000), Emu(10728960))

    # Usage criteria
    add_textbox(slide, Emu(731520), Emu(3750000), Emu(5000000), Emu(365760),
                "使い分けの基準", font=FONT_JP, size=Emu(228600), bold=True, color=WHITE)

    criteria = [
        ("1〜3ファイルの作業", "→ スキル", ACCENT),
        ("数十ファイルの探索", "→ サブエージェント", ACCENT2),
        ("結果を会話で直接使いたい", "→ スキル", ACCENT),
        ("メインのコンテキストを守りたい", "→ サブエージェント", ACCENT2),
    ]
    cy = Emu(4200000)
    for task, result, color in criteria:
        add_textbox(slide, Emu(914400), cy, Emu(5000000), Emu(274320),
                    f"•  {task}", font=FONT_JP, size=Emu(177800), color=BODY)
        add_textbox(slide, Emu(5800000), cy, Emu(5000000), Emu(274320),
                    result, font=FONT_JP, size=Emu(177800), bold=True, color=color)
        cy += Emu(340000)

    add_footer(slide, "→ スキル = 手順書、サブエージェント = 専門チームメンバー",
               y=Emu(5800000))
    add_page_num(slide, 12)

    add_speaker_notes(slide,
        "PART 1の締めくくりとして、スキルとサブエージェントの使い分けを整理します。\n\n"
        "どちらもClaudeの能力を拡張しますが、動作が違います。\n\n"
        "スキルはメイン会話の中で動くので、結果がそのまま残ります。"
        "1〜3ファイル程度の軽量な作業に向いています。\n\n"
        "サブエージェントは独立コンテキストで動いて、要約だけ返ります。"
        "数十ファイルの探索のような重い作業に向いています。\n\n"
        "覚え方としては、スキルは「手順書」、サブエージェントは「専門チームメンバー」。\n"
        "手順書は自分で見て実行するもの。チームメンバーには仕事を任せて報告だけもらう。"
        "この違いを意識しておくと、使い分けに迷わなくなるかなと思います。\n\n"
        "これでPART 1は以上です。なんとなく、Claudeの能力を拡張する2つの仕組みがつかめたかなと思います。")


def build_hooks_slide(prs):
    """Slide 14: Hooks."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 2 — Hooks & MCP")
    add_slide_title(slide, "Hooksとは")

    add_textbox(slide, Emu(731520), Emu(1100000), Emu(10000000), Emu(365760),
                "Claude Codeの特定のタイミングで、自動で処理を実行する仕組み",
                font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    add_textbox(slide, Emu(731520), Emu(1500000), Emu(10000000), Emu(274320),
                "設定場所: .claude/settings.json",
                font=FONT_CODE, size=Emu(165100), color=MUTED)

    # 3 JSON examples side by side
    ex_w = Emu(3429000)
    gap = Emu(228600)
    ex_y = Emu(1900000)
    ex_h = Emu(2400000)

    examples = [
        ("例1: 編集後にlint自動修正", "PostToolUse",
         '{\n  "matcher": "Edit|Write",\n  "hooks": [{\n    "type": "command",\n    "command":\n      "npm run lint:fix"\n  }]\n}'),
        ("例2: テストを自動実行", "PostToolUse",
         '{\n  "matcher": "Edit|Write",\n  "hooks": [{\n    "type": "command",\n    "command":\n      "npm test -- --bail"\n  }]\n}'),
        ("例3: 完了通知", "Notification",
         '{\n  "hooks": [{\n    "type": "command",\n    "command":\n      "say \'完了しました\'"\n  }]\n}'),
    ]

    for i, (title, hook_type, code) in enumerate(examples):
        x = Emu(731520) + i * (ex_w + gap)
        add_textbox(slide, x, ex_y, ex_w, Emu(274320),
                    title, font=FONT_JP, size=Emu(152400), bold=True, color=WHITE)
        add_textbox(slide, x, ex_y + Emu(260000), ex_w, Emu(200000),
                    hook_type, font=FONT_CODE, size=Emu(127000), color=MUTED)
        add_code_block(slide, x, ex_y + Emu(500000), ex_w, ex_h - Emu(500000), code)

    add_footer(slide,
               "→ Hooksが失敗 → Claudeに通知 → Claudeが自動で修正を試みる",
               y=Emu(4600000))
    add_page_num(slide, 14)

    add_speaker_notes(slide,
        "Hooksは、Claude Codeの特定のタイミングで自動で処理を実行する仕組みです。\n\n"
        "ハンドラには4タイプあるんですが、一番よく使うのはcommand、つまりシェルコマンドの実行です。"
        "今日はこのcommandを中心に紹介します。ちなみに他にはhttp、prompt、agentというタイプもありますが、"
        "まずはcommandだけ覚えておけば十分です。\n\n"
        "3つ実践例を紹介します。例1と例2は同じタイミングの例なので、実際にはプロジェクトに合わせて"
        "どちらか選んでください。\n\n"
        "1つ目、ファイル編集後にlintを自動修正。Claudeがファイルを編集するたびにlintが走ります。\n\n"
        "2つ目、ファイル編集後にテストを自動実行。テストが失敗するとClaudeに通知されて、"
        "Claudeがテスト結果を見て自動で修正を試みてくれます。\n\n"
        "3つ目、完了通知。macOSの say コマンドで音声通知です。長い作業を投げておいて、"
        "終わったら教えてもらう、という使い方ですね。\n\n"
        "ここで大事なのは、Hooksが失敗するとClaudeにその結果が通知されて、Claudeが自分で修正を試みてくれること。"
        "自動チェックと自己修正の組み合わせが結構強力です。")


def build_mcp_simple_slide(prs):
    """Slide 15: MCP simplified."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 2 — Hooks & MCP")
    add_slide_title(slide, "MCPとは")

    add_textbox(slide, Emu(731520), Emu(1200000), Emu(10000000), Emu(365760),
                "MCP = Model Context Protocol — 外部サービスとの連携プロトコル",
                font=FONT_EN, size=Emu(203200), bold=True, color=ACCENT)

    # Connection diagram
    connections = [
        ("Figma", "デザインからコンポーネント生成", ACCENT),
        ("Chrome", "ブラウザ操作・E2Eデバッグ", ACCENT2),
        ("Jira", "チケット管理の自動化", ACCENT3),
    ]

    y = Emu(1700000)
    for service, desc, color in connections:
        add_rounded_rect(slide, Emu(731520), y, Emu(2200000), Emu(365760), fill=CARD_BG)
        add_textbox(slide, Emu(831520), y + Emu(68580), Emu(2000000), Emu(274320),
                    "Claude Code", font=FONT_EN, size=Emu(152400), bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)

        add_textbox(slide, Emu(3000000), y + Emu(22860), Emu(1200000), Emu(320000),
                    "← MCP →", font=FONT_EN, size=Emu(139700), color=color,
                    align=PP_ALIGN.CENTER)

        add_rounded_rect(slide, Emu(4200000), y, Emu(1600000), Emu(365760), fill=CARD_BG)
        add_textbox(slide, Emu(4300000), y + Emu(68580), Emu(1400000), Emu(274320),
                    service, font=FONT_EN, size=Emu(152400), bold=True, color=color,
                    align=PP_ALIGN.CENTER)

        add_textbox(slide, Emu(6100000), y + Emu(68580), Emu(5000000), Emu(274320),
                    desc, font=FONT_JP, size=Emu(177800), color=BODY)

        y += Emu(457200)

    # Judgment table
    add_divider(slide, Emu(731520), Emu(3300000), Emu(10728960))

    add_textbox(slide, Emu(731520), Emu(3450000), Emu(5000000), Emu(365760),
                "判断基準", font=FONT_JP, size=Emu(228600), bold=True, color=WHITE)

    col_w = [Emu(5000000), Emu(5728960)]
    col_x = [Emu(731520), Emu(731520 + 5000000)]

    hy = Emu(3900000)
    for hdr, cx, cw in zip(["状況", "推奨"], col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(365760), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), hy + Emu(68580), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(177800), bold=True, color=ACCENT)

    criteria = [
        ("CLIで十分（GitHub、AWS等）", "CLI使用（gh、aws等）"),
        ("複雑な操作の連携が必要", "MCP検討"),
        ("ブラウザ操作が必要", "MCP使用（Chrome DevTools）"),
    ]
    for i, (situation, recommendation) in enumerate(criteria):
        ry = hy + Emu(365760) + Emu(i * 370000)
        for j, (cell, cx, cw) in enumerate(zip([situation, recommendation], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(365760), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fc = WHITE if j == 0 else BODY
            add_textbox(slide, cx + Emu(91440), ry + Emu(68580), cw - Emu(182880), Emu(274320),
                        cell, font=FONT_JP, size=Emu(177800), color=fc)

    add_footer(slide, "→ 原則: CLIで足りるならCLI。MCPは「CLIで足りないとき」の選択肢。設定は .mcp.json")
    add_page_num(slide, 15)

    add_speaker_notes(slide,
        "MCPは外部サービスとの連携プロトコルです。Figmaやブラウザ、Jiraなどとの深い連携が必要な場合に使います。\n\n"
        "じゃあどういうときにMCPを使うのか。判断基準は基本的に、CLIで十分ならCLIを使う、ということです。"
        "Anthropic公式も「ghコマンドはGitHub MCPより多くのケースで適している」と言っていたりします。\n\n"
        "MCPは「CLIで足りないとき」の選択肢、と覚えておくのがいいかなと思います。"
        "設定方法は .mcp.json をリポジトリに置くだけです。")


def build_worktree_slide(prs):
    """Slide 17: worktree + ISSUE.md."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_part_label(slide, "PART 3 — worktree")
    add_slide_title(slide, "worktree + ISSUE.md")

    add_textbox(slide, Emu(731520), Emu(1100000), Emu(10000000), Emu(365760),
                "Claude Codeにはworktree機能が組み込まれている",
                font=FONT_JP, size=Emu(203200), bold=True, color=ACCENT)

    # Code block: worktree commands
    add_code_block(slide, Emu(731520), Emu(1500000), Emu(5200000), Emu(600000),
                   "claude --worktree feature-auth    # 自動でworktree作成\n"
                   "claude --worktree bugfix-payment  # 別ターミナルで並列実行")

    add_textbox(slide, Emu(731520), Emu(2200000), Emu(5200000), Emu(274320),
                "各worktreeは独立したファイル状態。設定もすべて引き継がれる",
                font=FONT_JP, size=Emu(165100), color=MUTED)

    # issue-worker tree
    add_textbox(slide, Emu(731520), Emu(2600000), Emu(10000000), Emu(365760),
                "/issue-worker スキルは --worktree にISSUE.md生成を足す",
                font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT2)

    add_code_block(slide, Emu(731520), Emu(3000000), Emu(5200000), Emu(1100000),
                   "/issue-worker 123 456 789\n"
                   "  ├── worktree A: issue-123\n"
                   "  │   └── ISSUE.md自動生成 → Claude起動\n"
                   "  ├── worktree B: issue-456\n"
                   "  │   └── ISSUE.md自動生成 → Claude起動\n"
                   "  └── worktree C: issue-789\n"
                   "      └── ISSUE.md自動生成 → Claude起動")

    # Right side: steps and builtin skills
    add_textbox(slide, Emu(6200000), Emu(1100000), Emu(5000000), Emu(365760),
                "スキルがworktree作成〜Claude起動まで自動化",
                font=FONT_JP, size=Emu(190500), bold=True, color=WHITE)

    steps = [
        ("1", "Issue情報を取得"),
        ("2", "worktree + ブランチを作成"),
        ("3", "ISSUE.mdを自動生成"),
        ("4", "新しいターミナルでClaude Code起動"),
    ]
    sy = Emu(1500000)
    for num, label in steps:
        num_s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(6200000), sy,
                                        Emu(320000), Emu(320000))
        num_s.line.fill.background()
        num_s.fill.solid()
        num_s.fill.fore_color.rgb = ACCENT
        tf = num_s.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        _set_font(run, FONT_EN, Emu(152400), True, WHITE)

        add_textbox(slide, Emu(6600000), sy + Emu(22860), Emu(4800000), Emu(320000),
                    label, font=FONT_JP, size=Emu(165100), color=BODY)
        sy += Emu(370000)

    # Built-in skills table
    add_textbox(slide, Emu(6200000), Emu(3100000), Emu(5000000), Emu(365760),
                "組み込みスキルもworktreeを活用",
                font=FONT_JP, size=Emu(190500), bold=True, color=ACCENT3)

    scol_w = [Emu(1800000), Emu(3660960)]
    scol_x = [Emu(6200000), Emu(8000000)]
    shy = Emu(3500000)

    for hdr, cx, cw in zip(["スキル", "用途"], scol_x, scol_w):
        add_shape(slide, cx, shy, cw, Emu(310000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), shy + Emu(55000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT3)

    wskills = [
        ("/batch", "大規模変更をworktreeで並列実行→PR"),
        ("/simplify", "3エージェントが並列で品質検証"),
    ]
    for i, (sk, usage) in enumerate(wskills):
        ry = shy + Emu(310000) + Emu(i * 300000)
        for j, (cell, cx, cw) in enumerate(zip([sk, usage], scol_x, scol_w)):
            add_shape(slide, cx, ry, cw, Emu(290000), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fn = FONT_CODE if j == 0 else FONT_JP
            fc = ACCENT3 if j == 0 else BODY
            add_textbox(slide, cx + Emu(68580), ry + Emu(40000), cw - Emu(137160), Emu(274320),
                        cell, font=fn, size=Emu(127000), color=fc)

    add_footer(slide, "→ 各worktreeにはCLAUDE.md・Skills・Hooks・サブエージェントの設定がすべて引き継がれる",
               y=Emu(4500000))
    add_page_num(slide, 17)

    add_speaker_notes(slide,
        "Claude Codeには --worktree フラグが組み込まれています。"
        "claude --worktree feature-auth のように起動するだけで、自動的にworktreeが作成されて、"
        "独立した環境でセッションが始まります。\n\n"
        "--worktree だけでも並列作業はできるんですが、Issue番号との紐付けやISSUE.mdの自動生成までは"
        "やってくれません。そこで、PART 1で紹介した /issue-worker の出番です。"
        "Issue番号を渡すだけでworktree作成、ISSUE.md生成、Claude起動までが一気に進みます。\n\n"
        "先ほど紹介した /batch と /simplify は、実はworktreeを活用しています。"
        "/batch は大規模な変更を自動で分解し、各ユニットをworktree隔離で並列実行してPRまで作ってくれます。"
        "たとえば「src/以下のログ出力をすべてstructured loggingに移行して」と指示すると、"
        "対象を調査して5〜30のユニットに分けて、それぞれ独立して実行します。\n\n"
        "/simplify は実装後のコードレビュー用です。3つのレビューエージェントが並列で品質・効率・再利用性を"
        "チェックし、問題があれば自動修正してくれます。\n\n"
        "という感じで、/issue-worker のようなカスタムスキルと、/batch /simplify のような組み込みスキルを"
        "組み合わせることで、開発フロー全体をカバーできるようになります。\n\n"
        "ここで大事なポイントなんですが、--worktree で作られた環境にはCLAUDE.md、Skills、Hooks、"
        "サブエージェントの設定がすべて引き継がれます。つまり、PART 1とPART 2で整えた仕組みが、"
        "そのまま各並列作業で活きてくるわけです。")


def build_overview_slide(prs):
    """Slide 18: Big picture."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_slide_title(slide, "大規模開発の全体像", y=Emu(365760))

    add_textbox(slide, Emu(731520), Emu(914400), Emu(10000000), Emu(365760),
                "Claude Codeを強化する仕組みの全体像",
                font=FONT_JP, size=Emu(203200), bold=True, color=MUTED)

    # 6-row table
    col_w = [Emu(2400000), Emu(3600000), Emu(4728960)]
    col_x = [Emu(731520)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    hy = Emu(1300000)
    rh = Emu(310000)

    for hdr, cx, cw in zip(["仕組み", "役割", "設定場所"], col_x, col_w):
        add_shape(slide, cx, hy, cw, Emu(340000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(91440), hy + Emu(60000), cw - Emu(182880), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(165100), bold=True, color=ACCENT)

    items = [
        ("CLAUDE.md", "常識を教える", "プロジェクトルート"),
        ("Skills", "手順を仕組み化", ".claude/skills/"),
        ("Hooks", "品質チェックを自動化", ".claude/settings.json"),
        ("サブエージェント", "重い作業を委譲", ".claude/agents/"),
        ("MCP", "外部サービスとつなぐ", ".mcp.json"),
        ("worktree", "並列開発を実現", "Git機能"),
    ]
    colors = [ACCENT, ACCENT, ACCENT2, ACCENT, ACCENT3, ACCENT3]
    for i, ((name, role, loc), clr) in enumerate(zip(items, colors)):
        ry = hy + Emu(340000) + Emu(i * rh)
        for j, (cell, cx, cw) in enumerate(zip([name, role, loc], col_x, col_w)):
            add_shape(slide, cx, ry, cw, rh, fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            if j == 0:
                fn, fc, fb = FONT_EN, clr, True
            elif j == 2:
                fn, fc, fb = FONT_CODE, MUTED, False
            else:
                fn, fc, fb = FONT_JP, BODY, False
            add_textbox(slide, cx + Emu(91440), ry + Emu(50000), cw - Emu(182880), Emu(274320),
                        cell, font=fn, size=Emu(152400), bold=fb, color=fc)

    # 5-step flow
    add_divider(slide, Emu(731520), Emu(3500000), Emu(10728960))

    add_textbox(slide, Emu(731520), Emu(3600000), Emu(5000000), Emu(365760),
                "大規模開発のフロー", font=FONT_JP, size=Emu(228600), bold=True, color=WHITE)

    flow_steps = [
        ("1", "セットアップ", "CLAUDE.md / Skills / Hooks /\nサブエージェント / MCP を設定"),
        ("2", "Issue割り当て", "/issue-worker で\nworktree + ISSUE.md を作成"),
        ("3", "並列実行", "各worktreeでClaude Codeを起動し\n独立して作業を進める"),
        ("4", "品質担保", "Hooksが自動でlint・テスト"),
        ("5", "完了", "PR作成 → レビュー →\nマージ → worktree削除"),
    ]

    card_w = Emu(2057000)
    start_x = Emu(731520)
    flow_y = Emu(4050000)
    flow_h = Emu(1800000)
    gap = Emu(114300)

    for i, (num, title, desc) in enumerate(flow_steps):
        cx = start_x + i * (card_w + gap)
        add_rounded_rect(slide, cx, flow_y, card_w, flow_h, fill=CARD_BG)
        add_shape(slide, cx, flow_y, card_w, Emu(45720), fill=ACCENT)

        num_s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Emu(91440), flow_y + Emu(137160),
                                        Emu(320000), Emu(320000))
        num_s.line.fill.background()
        num_s.fill.solid()
        num_s.fill.fore_color.rgb = ACCENT
        tf = num_s.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        _set_font(run, FONT_EN, Emu(165100), True, WHITE)

        add_textbox(slide, cx + Emu(480000), flow_y + Emu(160000), card_w - Emu(548640), Emu(274320),
                    title, font=FONT_JP, size=Emu(165100), bold=True, color=WHITE)

        add_textbox(slide, cx + Emu(91440), flow_y + Emu(570000), card_w - Emu(182880), Emu(1100000),
                    desc, font=FONT_JP, size=Emu(127000), color=BODY)

    add_footer(slide, "→ すべてgit管理可能。チーム全員が同じ環境で並列開発できる",
               y=Emu(6100000))
    add_page_num(slide, 18)

    add_speaker_notes(slide,
        "ここまでの全体像をまとめます。\n\n"
        "CLAUDE.mdで常識を教え、Skillsで手順を仕組み化し、Hooksで品質チェックを自動化し、"
        "サブエージェントで重い作業を委譲し、MCPで外部サービスと連携する。\n"
        "そしてworktreeで、これらがすべて揃った環境を複数同時に立ち上げて並列開発する。"
        "全3回で学んだことがここに集約されている、という感じです。\n\n"
        "大規模開発のフローとしては、まずセットアップ。次にIssueごとにworktreeを作成。"
        "各worktreeでClaude Codeを起動し、独立して作業を進めます。"
        "Hooksが自動でチェックし、最後にPR作成。自動チェックはレビュー負荷を下げるためのもので、"
        "最終的なコードレビューとマージ判断は人間が行います。\n\n"
        "実際の運用実績もお伝えしますと、この仕組みを導入しているプロジェクトでは、"
        "worktreeを19個同時に稼働させて、週に3〜5個のIssueを並列で進めています。"
        "MRは1日1〜2件のペースで出ています。手作業で5〜10分かかっていた環境構築が、"
        "worktreeの自動生成で1分以内に終わるようになりました。")


def build_demo_slide(prs, title, slide_num):
    """Demo slide (minimal text)."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_textbox(slide, Emu(0), Emu(2200000), Emu(SLIDE_W), Emu(762000),
                "DEMO", font=FONT_EN, size=Emu(762000), bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    add_textbox(slide, Emu(0), Emu(3200000), Emu(SLIDE_W), Emu(457200),
                title, font=FONT_JP, size=Emu(254000),
                color=MUTED, align=PP_ALIGN.CENTER)

    add_page_num(slide, slide_num)

    add_speaker_notes(slide,
        "では実際にやってみましょう。\n\n"
        "（/issue-worker を実行して）worktreeが作成されて、ISSUE.mdが自動生成されます。\n"
        "このworktreeには、CLAUDE.mdもSkillsもHooksの設定もすべて引き継がれています。\n\n"
        "Claude Codeを起動して「ISSUE.mdを読んで計画を立てて」と指示します。\n"
        "サブエージェントが関連コードを調査し、計画を提案してくれます。\n"
        "「OK、実装して」と言えば実装が始まり、Hooksによるlintチェックが自動で走ります。\n\n"
        "で、別のターミナルでは別のworktreeで別のIssueが並行して進んでいる。"
        "これが今日お話しした大規模開発の全体像です。")


def build_summary_slide(prs):
    """Slide 20: Summary."""
    slide = new_slide(prs)
    set_bg(slide, BG_MAIN)

    add_slide_title(slide, "まとめ / シリーズ総括", y=Emu(365760))

    add_textbox(slide, Emu(731520), Emu(914400), Emu(10000000), Emu(365760),
                "3つの持ち帰り", font=FONT_JP, size=Emu(254000), bold=True, color=MUTED)

    takeaways = [
        ("1", "Skills & サブエージェントで\n指示と作業を仕組み化する",
         "スキル = 手順書（メイン会話内）\nサブエージェント = 専門チーム\nメンバー（独立コンテキスト）", ACCENT),
        ("2", "Hooks & MCPで\n品質と連携を自動化する",
         "品質チェックの自動実行、\n外部サービスとの連携。\nCLIで足りるならCLI", ACCENT2),
        ("3", "worktreeで\nすべてを束ねて並列開発する",
         "設定はすべてのworktreeに\n引き継がれる。Issue番号を\n渡すだけで作業環境が整う", ACCENT3),
    ]

    card_w = Emu(3429000)
    card_h = Emu(2743200)
    start_x = Emu(731520)
    start_y = Emu(1371600)
    gap = Emu(228600)

    for i, (num, title, desc, color) in enumerate(takeaways):
        cx = start_x + i * (card_w + gap)

        add_rounded_rect(slide, cx, start_y, card_w, card_h, fill=CARD_BG)
        add_shape(slide, cx, start_y, card_w, Emu(54864), fill=color)

        num_s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Emu(182880), start_y + Emu(182880),
                                        Emu(548640), Emu(548640))
        num_s.line.fill.background()
        num_s.fill.solid()
        num_s.fill.fore_color.rgb = color
        tf = num_s.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        _set_font(run, FONT_EN, Emu(254000), True, WHITE)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(822960), card_w - Emu(365760), Emu(640080),
                    title, font=FONT_JP, size=Emu(190500), bold=True, color=WHITE)

        add_textbox(slide, cx + Emu(182880), start_y + Emu(1554480), card_w - Emu(365760), Emu(1097280),
                    desc, font=FONT_JP, size=Emu(177800), color=BODY)

    # Series review table
    add_textbox(slide, Emu(731520), Emu(4400000), Emu(10728960), Emu(365760),
                "全3回の振り返り", font=FONT_JP, size=Emu(203200), bold=True, color=WHITE)

    col_w = [Emu(1200000), Emu(2400000), Emu(7128960)]
    col_x = [Emu(731520)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    rhy = Emu(4750000)
    for hdr, cx, cw in zip(["回", "テーマ", "学んだこと"], col_x, col_w):
        add_shape(slide, cx, rhy, cw, Emu(300000), fill=CARD_BG2)
        add_textbox(slide, cx + Emu(68580), rhy + Emu(45000), cw - Emu(137160), Emu(274320),
                    hdr, font=FONT_JP, size=Emu(139700), bold=True, color=ACCENT)

    review_rows = [
        ("第1回", "理解する", "基本操作 / コミュニケーション / Plan Mode"),
        ("第2回", "定着させる", "CLAUDE.md / settings.json"),
        ("第3回", "大規模開発に挑む", "Skills / サブエージェント / Hooks / MCP / worktree"),
    ]
    for i, (session, theme, learned) in enumerate(review_rows):
        ry = rhy + Emu(300000) + Emu(i * 280000)
        for j, (cell, cx, cw) in enumerate(zip([session, theme, learned], col_x, col_w)):
            add_shape(slide, cx, ry, cw, Emu(270000), fill=CARD_BG if i % 2 == 0 else CARD_BG3)
            fc = WHITE if j <= 1 else BODY
            fb = j <= 1
            add_textbox(slide, cx + Emu(68580), ry + Emu(35000), cw - Emu(137160), Emu(274320),
                        cell, font=FONT_JP, size=Emu(127000), bold=fb, color=fc)

    add_divider(slide, Emu(731520), Emu(6000000), Emu(10728960))

    add_textbox(slide, Emu(731520), Emu(6100000), Emu(10728960), Emu(365760),
                "まず使ってみて、プロジェクトに合わせて育てていってください",
                font=FONT_JP, size=Emu(190500), color=MUTED)

    add_page_num(slide, 20)

    add_speaker_notes(slide,
        "まとめです。今日の持ち帰りは3つ。\n\n"
        "1つ目。SkillsとサブエージェントでClaudeの能力を拡張する。"
        "軽い作業はスキル、重い探索はサブエージェント。手順書とチームメンバーの違いでしたね。\n\n"
        "2つ目。HooksとMCPで品質と連携を自動化する。Hooksは編集のたびに自動でチェック、"
        "MCPはCLIでは足りない外部連携に。\n\n"
        "3つ目。worktreeですべてを束ねて並列開発する。セットアップした環境がそのまま各worktreeに引き継がれます。\n\n"
        "全3回を通して、「理解する、定着させる、大規模開発に挑む」という流れでお話ししてきました。\n\n"
        "全3回にお付き合いいただき、ありがとうございました。\n"
        "まずは使ってみて、プロジェクトに合わせて育てていっていただければと思います。\n"
        "では、質疑応答に移ります。")


# ── Main ───────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)                                          # 01
    build_roadmap_slide(prs)                                        # 02
    build_agenda_slide(prs)                                         # 03
    build_section_slide(prs, 1, "Skills & サブエージェント",
                        "繰り返しをなくし、分業する", 4)             # 04
    build_skill_intro_slide(prs)                                    # 05
    build_skill_types_slide(prs)                                    # 06
    build_skill_examples_slide(prs)                                 # 07
    build_skill_design_slide(prs)                                   # 08
    build_builtin_commands_slide(prs)                               # 09
    build_subagent_intro_slide(prs)                                 # 10
    build_subagent_combined_slide(prs)                              # 11
    build_skill_vs_subagent_slide(prs)                              # 12
    build_section_slide(prs, 2, "Hooks & MCP", "自動化し、つなぐ", 13)  # 13
    build_hooks_slide(prs)                                          # 14
    build_mcp_simple_slide(prs)                                     # 15
    build_section_slide(prs, 3, "worktree",
                        "すべてを束ねて、並列で挑む", 16)           # 16
    build_worktree_slide(prs)                                       # 17
    build_overview_slide(prs)                                       # 18
    build_demo_slide(prs, "worktree + ISSUE.md", 19)                # 19
    build_summary_slide(prs)                                        # 20

    prs.save("slides/Claude_session3.pptx")
    print(f"✓ slides/Claude_session3.pptx を生成しました（全{TOTAL_SLIDES}スライド）")


if __name__ == "__main__":
    main()
